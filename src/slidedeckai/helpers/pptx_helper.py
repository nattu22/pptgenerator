"""
A set of functions to create a PowerPoint slide deck.
"""
import logging
import os
import pathlib
import random
import re
import tempfile
from typing import Optional

import json5
import pptx
from dotenv import load_dotenv
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.shapes.placeholder import PicturePlaceholder, SlidePlaceholder

from ..layout_analyzer import TemplateAnalyzer
from ..content_matcher import ContentLayoutMatcher

from . import icons_embeddings as ice
from . import image_search as ims
from ..global_config import GlobalConfig
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
load_dotenv()

# English Metric Unit (used by PowerPoint) to inches
EMU_TO_INCH_SCALING_FACTOR = 1.0 / 914400
INCHES_3 = pptx.util.Inches(3)
INCHES_2 = pptx.util.Inches(2)
INCHES_1_5 = pptx.util.Inches(1.5)
INCHES_1 = pptx.util.Inches(1)
INCHES_0_8 = pptx.util.Inches(0.8)
INCHES_0_9 = pptx.util.Inches(0.9)
INCHES_0_5 = pptx.util.Inches(0.5)
INCHES_0_4 = pptx.util.Inches(0.4)
INCHES_0_3 = pptx.util.Inches(0.3)
INCHES_0_2 = pptx.util.Inches(0.2)

STEP_BY_STEP_PROCESS_MARKER = '>> '
ICON_BEGINNING_MARKER = '[['
ICON_END_MARKER = ']]'

ICON_SIZE = INCHES_0_8
ICON_BG_SIZE = INCHES_1

IMAGE_DISPLAY_PROBABILITY = 1 / 3.0
FOREGROUND_IMAGE_PROBABILITY = 0.8

SLIDE_NUMBER_REGEX = re.compile(r"^slide[ ]+\d+:", re.IGNORECASE)
ICONS_REGEX = re.compile(r"\[\[(.*?)\]\]\s*(.*)")
BOLD_ITALICS_PATTERN = re.compile(r'(\*\*(.*?)\*\*|\*(.*?)\*)')

ICON_COLORS = [
    pptx.dml.color.RGBColor.from_string('800000'),  # Maroon
    pptx.dml.color.RGBColor.from_string('6A5ACD'),  # SlateBlue
    pptx.dml.color.RGBColor.from_string('556B2F'),  # DarkOliveGreen
    pptx.dml.color.RGBColor.from_string('2F4F4F'),  # DarkSlateGray
    pptx.dml.color.RGBColor.from_string('4682B4'),  # SteelBlue
    pptx.dml.color.RGBColor.from_string('5F9EA0'),  # CadetBlue
]


logger = logging.getLogger(__name__)
logging.getLogger('PIL.PngImagePlugin').setLevel(logging.ERROR)

def remove_unused_placeholders(slide):
    """Remove empty placeholders from a slide (keep title)."""
    shapes_to_remove = []
    
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.idx == 0:  # Keep title
            continue
            
        is_empty = True
        if shape.has_text_frame and shape.text_frame.text.strip():
            is_empty = False
        
        if is_empty:
            shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        try:
            sp = shape.element
            sp.getparent().remove(sp)
        except Exception as e:
            logger.warning(f"Could not remove placeholder: {e}")

def get_content_placeholders_left_to_right(slide):
    """Get content placeholders ordered from left to right (skip title AND subtitle)."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    
    content_placeholders = []
    
    for shape in slide.placeholders:
        # Skip title (idx 0)
        if shape.placeholder_format.idx == 0:
            continue
        
        # Skip subtitle (idx 1)
        if shape.placeholder_format.idx == 1:
            continue
            
        try:
            ph_type = shape.placeholder_format.type
            # Skip TITLE and SUBTITLE types
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.SUBTITLE):
                continue
        except:
            pass
        
        content_placeholders.append((shape, shape.left))
    
    if not content_placeholders:
        return []
    
    content_placeholders.sort(key=lambda x: x[1])
    return [ph[0] for ph in content_placeholders]

def get_largest_content_placeholder(slide):
    """
    Find the largest non-title, non-subtitle placeholder for content.
    """
    from pptx.enum.shapes import PP_PLACEHOLDER
    
    content_placeholders = []
    
    for shape in slide.placeholders:
        # Skip title (idx 0)
        if shape.placeholder_format.idx == 0:
            continue
        
        # Skip subtitle (idx 1)
        if shape.placeholder_format.idx == 1:
            continue
            
        try:
            ph_type = shape.placeholder_format.type
            # Skip TITLE and SUBTITLE placeholders
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.SUBTITLE):
                continue
        except:
            pass
            
        # Calculate area
        area = shape.width * shape.height
        content_placeholders.append((shape, area))
    
    if not content_placeholders:
        return None
    
    # Return the placeholder with the largest area
    return max(content_placeholders, key=lambda x: x[1])[0]
    
def get_placeholder_safely(slide, idx, placeholder_name=""):
    """
    Safely get a placeholder from a slide by index.
    Falls back to the largest non-title placeholder if index not found.
    
    Args:
        slide: The slide object
        idx: The placeholder index to retrieve
        placeholder_name: Optional name for better error messages
    
    Returns:
        The placeholder object if found, None otherwise
    """
    try:
        return slide.placeholders[idx]
    except KeyError:
        # Get all available placeholder indices
        available_indices = [ph.placeholder_format.idx for ph in slide.placeholders]
        
        logger.warning(
            f"Placeholder {idx} ({placeholder_name}) not found in slide. "
            f"Available indices: {available_indices}"
        )
        
        # Try to find the largest non-title placeholder
        largest_placeholder = get_largest_content_placeholder(slide)
        
        if largest_placeholder:
            logger.info(f"Using largest content placeholder at index {largest_placeholder.placeholder_format.idx}")
            return largest_placeholder
        
        # Last resort: try first available non-title placeholder
        non_title_indices = [idx for idx in available_indices if idx != 0]
        if non_title_indices:
            fallback_idx = non_title_indices[0]
            logger.info(f"Using fallback placeholder at index {fallback_idx}")
            return slide.placeholders[fallback_idx]
        
        logger.error("No placeholders found in this slide at all!")
        return None


def get_placeholder_by_type(slide, placeholder_type):
    """
    Get placeholder by type instead of index (more robust).
    
    Args:
        slide: The slide object
        placeholder_type: Type of placeholder (e.g., PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.BODY)
    
    Returns:
        The first placeholder of the specified type, or None
    """
    for shape in slide.placeholders:
        if shape.placeholder_format.type == placeholder_type:
            return shape
    return None
    
def remove_slide_number_from_heading(header: str) -> str:
    """Remove slide number prefix like 'Slide 1:' from heading."""
    import re
    SLIDE_NUMBER_REGEX = re.compile(r"^slide[ ]+\d+:", re.IGNORECASE)
    if SLIDE_NUMBER_REGEX.match(header):
        idx = header.find(':')
        header = header[idx + 1:].strip()
    return header


def add_bulleted_items(text_frame, flat_items_list: list):
    """
    Add bullet points to a text frame with proper formatting.
    
    Args:
        text_frame: The text frame to add bullets to
        flat_items_list: List of (text, level) tuples
    """
    for idx, an_item in enumerate(flat_items_list):
        if idx == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
            paragraph.level = an_item[1]
        
        STEP_BY_STEP_PROCESS_MARKER = '>> '
        format_text(paragraph, an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER))

def format_text(frame_paragraph, text: str):
    """
    Apply bold (**text**) and italic (*text*) formatting.
    
    Args:
        frame_paragraph: The paragraph to format
        text: Text with markdown-style formatting
    """
    import re
    BOLD_ITALICS_PATTERN = re.compile(r'(\*\*(.*?)\*\*|\*(.*?)\*)')
    
    matches = list(BOLD_ITALICS_PATTERN.finditer(text))
    last_index = 0
    
    for match in matches:
        start, end = match.span()
        if start > last_index:
            run = frame_paragraph.add_run()
            run.text = text[last_index:start]
        
        if match.group(2):  # Bold
            run = frame_paragraph.add_run()
            run.text = match.group(2)
            run.font.bold = True
        elif match.group(3):  # Italics
            run = frame_paragraph.add_run()
            run.text = match.group(3)
            run.font.italic = True
        
        last_index = end
    
    if last_index < len(text):
        run = frame_paragraph.add_run()
        run.text = text[last_index:]

def generate_powerpoint_presentation(
        parsed_data: dict,
        slides_template: str,
        output_file_path: pathlib.Path
) -> list:
    """FULLY DYNAMIC WITH STORY AWARENESS"""
    
    presentation = pptx.Presentation(GlobalConfig.PPTX_TEMPLATE_FILES[slides_template]['file'])
    
    # ANALYZE TEMPLATE
    logger.info("🔍 Analyzing template structure...")
    analyzer = TemplateAnalyzer(presentation)
    analyzer.print_summary()
    
    # CREATE CONTENT MATCHER
    matcher = ContentLayoutMatcher(analyzer)
    
    # Clear existing slides
    slide_id_list = [slide.slide_id for slide in presentation.slides]
    for slide_id in slide_id_list:
        rId = presentation.slides._sldIdLst[0].rId
        presentation.part.drop_rel(rId)
        del presentation.slides._sldIdLst[0]
    
    slide_width_inch, slide_height_inch = _get_slide_width_height_inches(presentation)
    # Derive theme colors, fonts and responsive sizes from template (avoid hardcoded values)
    from pptx.util import Inches
    global ICON_SIZE, ICON_BG_SIZE, ICON_COLORS, THEME_COLORS, INCHES_1, INCHES_0_5, INCHES_2, INCHES_3, INCHES_0_2
    global THEME_FONT_NAME, THEME_FONT_SIZE

    # Default sizing as proportion of slide width
    ICON_SIZE = Inches(min(max(slide_width_inch * 0.12, 0.4), 1.2))
    ICON_BG_SIZE = Inches(ICON_SIZE.inches * 1.25)

    # Default small inch constants (used throughout helpers) - keep sensible defaults
    INCHES_1 = Inches(1)
    INCHES_0_5 = Inches(0.5)
    INCHES_2 = Inches(2)
    INCHES_3 = Inches(3)
    INCHES_0_2 = Inches(0.2)

    # Extract theme colors from slide master if available
    THEME_COLORS = {}
    try:
        if presentation.slide_master and getattr(presentation.slide_master, 'theme', None):
            theme = presentation.slide_master.theme
            # theme.theme_colors may be sequence-like; extract by index when available
            tc = getattr(theme, 'theme_colors', None)
            if tc:
                THEME_COLORS = {
                    'text1': tc[0] if len(tc) > 0 else None,
                    'text2': tc[1] if len(tc) > 1 else None,
                    'background1': tc[2] if len(tc) > 2 else None,
                    'accent1': tc[4] if len(tc) > 4 else None,
                    'accent2': tc[5] if len(tc) > 5 else None,
                }
    except Exception:
        THEME_COLORS = {}

    # Icon colors: prefer theme accent colors, otherwise keep existing palette
    try:
        if THEME_COLORS.get('accent1') or THEME_COLORS.get('accent2'):
            ICON_COLORS = [THEME_COLORS.get('accent1') or pptx.dml.color.RGBColor.from_string('4682B4'),
                           THEME_COLORS.get('accent2') or pptx.dml.color.RGBColor.from_string('5F9EA0')]
    except Exception:
        pass

    # Extract default font from first layout placeholder if possible
    THEME_FONT_NAME = 'Calibri'
    THEME_FONT_SIZE = Pt(18)
    try:
        for layout in presentation.slide_layouts:
            for shape in layout.placeholders:
                if getattr(shape, 'has_text_frame', False) and shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        run = para.runs[0]
                        THEME_FONT_NAME = run.font.name or THEME_FONT_NAME
                        THEME_FONT_SIZE = run.font.size or THEME_FONT_SIZE
                        raise StopIteration
    except StopIteration:
        pass
    except Exception:
        pass
    
    # TITLE SLIDE
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    presentation_title = parsed_data.get('title', 'Untitled Presentation')
    
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = presentation_title
    
    subtitle_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0:
            subtitle_placeholder = shape
            break
    
    if subtitle_placeholder:
        try:
            subtitle_placeholder.text = 'Generated by SlideDeck AI'
        except Exception as e:
            logger.warning(f"Could not set subtitle: {e}")
    
    all_headers = [presentation_title]
    
    # MAIN CONTENT LOOP - WITH STORY AWARENESS
    total_slides = len(parsed_data['slides'])
    used_content_types = []  # Track to enforce diversity
    
    for slide_index, a_slide in enumerate(parsed_data['slides']):
        try:
            logger.info(f"\n{'='*60}")
            logger.info(f"Processing slide {slide_index + 1}/{total_slides}: {a_slide.get('heading', 'Untitled')}")
            
            # DETECT CONTENT TYPE
            content_type = matcher._infer_content_type_from_json(a_slide)
            
            # ENFORCE DIVERSITY: No 3 consecutive same types
            if len(used_content_types) >= 2:
                if used_content_types[-1] == used_content_types[-2] == content_type:
                    logger.warning(f"⚠️ Avoiding 3rd consecutive {content_type} slide")
                    # Try to use alternative content if available
                    if 'bullet_points' in a_slide and content_type != 'bullets':
                        content_type = 'bullets'  # Fallback to bullets
            
            used_content_types.append(content_type)
            
            # SMART LAYOUT SELECTION WITH STORY AWARENESS
            layout_idx = matcher.select_layout_for_slide(
                a_slide, 
                slide_index=slide_index,
                total_slides=total_slides
            )
            
            # GET LAYOUT CAPABILITY
            layout_capability = analyzer.layouts.get(layout_idx)
            
            if not layout_capability:
                logger.error(f"Layout {layout_idx} not found!")
                continue
            
            # CREATE SLIDE
            slide_layout = presentation.slide_layouts[layout_idx]
            slide = presentation.slides.add_slide(slide_layout)
            
            # SET TITLE
            title_shape = slide.shapes.title
            if title_shape:
                heading = remove_slide_number_from_heading(a_slide['heading'])
                title_shape.text = heading
                all_headers.append(heading)
            
            # LOG LAYOUT SELECTION
            logger.info(f"✓ Using layout {layout_idx}: {layout_capability.name}")
            logger.info(f"  Story type: {layout_capability.semantic_story_type}")
            logger.info(f"  Exec score: {layout_capability.executive_suitability:.0f}/100")
            logger.info(f"  Content type: {content_type}")
            
            # GET CONTENT DENSITY RECOMMENDATION
            density_rec = layout_capability.content_density_recommendation
            if density_rec:
                logger.info(f"  Density: {density_rec.get('density_style', 'standard')}, "
                           f"{density_rec.get('total_words_target', 0)} words target")
            
            # ROUTE TO APPROPRIATE HANDLER
            if 'chart' in a_slide and a_slide['chart']:
                logger.info("  → Rendering chart")
                _handle_chart_dynamic(slide, a_slide, layout_capability, analyzer)
                
            elif 'table' in a_slide and a_slide['table']:
                logger.info("  → Rendering table")
                _handle_table_dynamic(slide, a_slide, layout_capability, analyzer)
                
            elif matcher._is_icon_slide(a_slide):
                logger.info("  → Rendering pictograms")
                _handle_icons_dynamic(slide, a_slide, layout_capability, analyzer, presentation)
                
            elif matcher._is_double_column_slide(a_slide):
                logger.info("  → Rendering double column")
                _handle_double_column_dynamic(slide, a_slide, layout_capability, analyzer)
                
            elif 'bullet_points' in a_slide:
                logger.info("  → Rendering bullets")
                # RESPECT DENSITY RECOMMENDATION
                if density_rec and density_rec.get('avoid_overflow'):
                    max_bullets = density_rec.get('bullets_recommended', 10)
                    logger.info(f"  Target bullets: {max_bullets}")
                
                _handle_bullets_dynamic(slide, a_slide, layout_capability, analyzer)
            
            else:
                logger.warning("  ⚠️ No content to render")
            
            # ADD KEY MESSAGE IF PRESENT
            if 'key_message' in a_slide and a_slide['key_message']:
                _handle_key_message(slide, a_slide, slide_width_inch, slide_height_inch)
            
            # REMOVE UNUSED PLACEHOLDERS
            remove_unused_placeholders(slide)
            
            logger.info(f"✓ Slide {slide_index + 1} completed")
            
        except Exception as e:
            logger.error(f'❌ Error processing slide {slide_index + 1}: {e}', exc_info=True)
            continue
    
    # THANK YOU SLIDE
    logger.info("\nAdding thank you slide...")
    last_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(last_slide_layout)
    
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = 'Thank you!'
    
    # SAVE
    presentation.save(output_file_path)
    
    logger.info(f"\n{'='*60}")
    logger.info(f"✅ Presentation saved: {output_file_path}")
    logger.info(f"   Total slides: {len(presentation.slides)}")
    logger.info(f"   Content diversity: {len(set(used_content_types))} types")
    logger.info(f"{'='*60}\n")
    
    return all_headers

def get_flat_list_of_contents(items: list, level: int) -> list:
    """
    Flatten hierarchical bullet points into (text, level) tuples.
    
    Args:
        items: List of strings or nested lists
        level: Current hierarchy level
    
    Returns:
        List of (text, level) tuples
    """
    flat_list = []
    for item in items:
        if isinstance(item, str):
            flat_list.append((item, level))
        elif isinstance(item, list):
            flat_list = flat_list + get_flat_list_of_contents(item, level + 1)
    return flat_list


def get_slide_placeholders(
        slide: pptx.slide.Slide,
        layout_number: int,
        is_debug: bool = False
) -> list[tuple[int, str]]:
    """
    Return the index and name (lower case) of all placeholders present in a
    slide, except the title placeholder.

    A placeholder in a slide is a place to add content. Each placeholder has a
    name and an index. This index is not a list index; it is a key used to look up
    a dict and may be non-contiguous. The title placeholder always has index 0.
    User-added placeholders get indices starting from 10.

    With user-edited or added placeholders, indices may be difficult to track. This
    function returns the placeholders' names as well, which may help distinguish
    between placeholders.

    Args:
        slide: The slide.
        layout_number: The layout number used by the slide.
        is_debug: Whether to print debugging statements.

    Returns:
        list[tuple[int, str]]: A list of (index, name) tuples for placeholders
        present in the slide, excluding the title placeholder.
    """

    if is_debug:
        print(
            f'Slide layout #{layout_number}:'
            f' # of placeholders: {len(slide.shapes.placeholders)} (including the title)'
        )

    placeholders = [
        (shape.placeholder_format.idx, shape.name.lower()) for shape in slide.shapes.placeholders
    ]
    placeholders.pop(0)  # Remove the title placeholder

    if is_debug:
        print(placeholders)

    return placeholders

def _add_text_at_bottom(
        slide: pptx.slide.Slide,
        slide_width_inch: float,
        slide_height_inch: float,
        text: str,
        hyperlink: Optional[str] = None,
        target_height: Optional[float] = 0.5
):
    """
    Add arbitrary text to a textbox positioned near the lower-left side of a slide.

    Args:
        slide: The slide.
        slide_width_inch: The width of the slide in inches.
        slide_height_inch: The height of the slide in inches.
        text: The text to be added.
        hyperlink: Optional; the hyperlink to be added to the text.
        target_height: Optional[float]; the target height of the box in inches.
    """

    footer = slide.shapes.add_textbox(
        left=INCHES_1,
        top=pptx.util.Inches(slide_height_inch - target_height),
        width=pptx.util.Inches(slide_width_inch),
        height=pptx.util.Inches(target_height)
    )

    paragraph = footer.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    # Use theme font size if available, otherwise fallback
    try:
        run.font.size = Pt(int(THEME_FONT_SIZE.pt * 0.6))
    except Exception:
        run.font.size = Pt(10)
    run.font.underline = False

    if hyperlink:
        run.hyperlink.address = hyperlink
        
def _handle_key_message(
        the_slide: pptx.slide.Slide,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
):
    """
    FIXED: Adaptive font color based on background brightness.
    
    Add a key message box at the bottom center of the slide.
    Text color automatically adapts:
    - Light backgrounds → Black text
    - Dark backgrounds → White text
    
    Args:
        the_slide: The slide object to add the key message to
        slide_json: Dictionary containing 'key_message' key
        slide_width_inch: Width of slide in inches
        slide_height_inch: Height of slide in inches
    """
    if 'key_message' not in slide_json or not slide_json['key_message']:
        return
    
    key_message_text = slide_json['key_message']
    
    # Calculate position and size (use template-driven sizes)
    try:
        height = Inches(min(max(slide_height_inch * 0.12, 0.8), 2.0))
        width = Inches(slide_width_inch / 2.3)
        top = Inches(slide_height_inch - height.inches - 0.1)
        left = Inches((slide_width_inch - width.inches) / 2)
    except Exception:
        height = Inches(1.6)
        width = Inches(slide_width_inch / 2.3)
        top = Inches(slide_height_inch - height.inches - 0.1)
        left = Inches((slide_width_inch - width.inches) / 2)
    
    # Create the rounded rectangle shape
    shape = the_slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left=left,
        top=top,
        width=width,
        height=height
    )
    
    # Configure text frame
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    # Add the key message text
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = pptx.enum.text.PP_ALIGN.CENTER
    format_text(paragraph, key_message_text)
    
    # Center the text vertically
    text_frame.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
    
    # Style the shape using theme background color when available
    shape.fill.solid()
    from pptx.dml.color import RGBColor
    try:
        bg_theme = THEME_COLORS.get('background1') if 'THEME_COLORS' in globals() else None
        if bg_theme is not None and hasattr(bg_theme, 'rgb'):
            bg_color = bg_theme.rgb
        elif isinstance(bg_theme, RGBColor):
            bg_color = bg_theme
        else:
            bg_color = RGBColor(240, 240, 240)
    except Exception:
        bg_color = RGBColor(240, 240, 240)

    try:
        shape.fill.fore_color.rgb = bg_color
    except Exception:
        try:
            shape.fill.fore_color.theme_color = pptx.enum.dml.MSO_THEME_COLOR.ACCENT_1
        except Exception:
            pass

    try:
        shape.line.color.rgb = RGBColor(200, 200, 200)
        shape.line.width = Pt(1)
    except Exception:
        pass

    # FIX: Determine text color based on background brightness
    # Using relative luminance formula (ITU-R BT.709)
    try:
        if hasattr(bg_color, 'rgb'):
            r, g, b = bg_color.rgb[0], bg_color.rgb[1], bg_color.rgb[2]
        elif isinstance(bg_color, RGBColor):
            r, g, b = bg_color.r, bg_color.g, bg_color.b
        else:
            r, g, b = 240, 240, 240
        luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    except Exception:
        luminance = 1.0
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    
    # Use black text for light backgrounds, white for dark
    if luminance > 0.5:
        text_color = pptx.dml.color.RGBColor(0, 0, 0)  # Black
    else:
        text_color = pptx.dml.color.RGBColor(255, 255, 255)  # White
    
    # Apply text formatting
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.font.size:
                run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = text_color


def _get_slide_width_height_inches(presentation: pptx.Presentation) -> tuple[float, float]:
    """
    Get the dimensions of a slide in inches.

    Args:
        presentation: The presentation object.

    Returns:
        The width and the height.
    """

    slide_width_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_width
    slide_height_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_height

    return slide_width_inch, slide_height_inch
    
# --------------------------------------------------------------------------------------------- #

def _handle_chart_dynamic(slide, slide_json: dict, layout_capability, analyzer):
    """Handle charts - use LARGEST content placeholder automatically"""
    if 'chart' not in slide_json or not slide_json['chart']:
        return
        
    chart_data_json = slide_json['chart']
    
    # FIX: Use get_largest_content_placeholder - it already skips title AND subtitle!
    actual_placeholder = get_largest_content_placeholder(slide)
    
    if not actual_placeholder:
        logger.error("No content placeholders for chart!")
        return
    
    # Use placeholder dimensions
    x = actual_placeholder.left
    y = actual_placeholder.top
    cx = actual_placeholder.width
    cy = actual_placeholder.height
    
    # Remove placeholder
    try:
        sp = actual_placeholder.element
        sp.getparent().remove(sp)
    except Exception as e:
        logger.warning(f"Could not remove placeholder: {e}")
    
    # Chart type mapping
    chart_type_map = {
        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'line': XL_CHART_TYPE.LINE,
        'pie': XL_CHART_TYPE.PIE,
    }
    
    chart_type = chart_data_json.get('type', 'column').lower()
    xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # Prepare data
    chart_data = CategoryChartData()
    chart_data.categories = chart_data_json.get('categories', [])
    
    series_data = chart_data_json.get('series', [])
    
    if chart_type == 'pie':
        if isinstance(series_data[0], dict):
            chart_data.add_series(series_data[0].get('name', 'Values'), series_data[0].get('values', []))
        else:
            chart_data.add_series('Values', series_data)
    else:
        for series_item in series_data:
            if isinstance(series_item, dict):
                chart_data.add_series(series_item.get('name', 'Series'), series_item.get('values', []))
            else:
                chart_data.add_series('Series', series_item)
    
    try:
        chart = slide.shapes.add_chart(xl_chart_type, x, y, cx, cy, chart_data).chart
        chart.has_legend = True
        
        if chart_type != 'pie':
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        
        if 'title' in chart_data_json:
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_data_json['title']
        # Apply theme font to chart title and legend where possible
        try:
            title_tf = chart.chart_title.text_frame
            for para in title_tf.paragraphs:
                for run in para.runs:
                    run.font.name = THEME_FONT_NAME
                    run.font.size = THEME_FONT_SIZE
        except Exception:
            pass

        try:
            if chart.has_legend and chart.legend and chart.legend.text_frame:
                for para in chart.legend.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = THEME_FONT_NAME
                        run.font.size = Pt(max(10, THEME_FONT_SIZE.pt * 0.8))
        except Exception:
            pass

        logger.info(f"✓ Chart added using largest placeholder")
    except Exception as e:
        logger.error(f"Failed to add chart: {e}")

def _handle_table_dynamic(slide, slide_json: dict, layout_capability, analyzer):
    """Handle tables - use LARGEST content placeholder automatically"""
    if 'table' not in slide_json or not slide_json['table']:
        return
        
    table_data = slide_json['table']
    headers = table_data.get('headers', [])
    rows = table_data.get('rows', [])
    
    if not headers or not rows:
        return
    
    # FIX: Use get_largest_content_placeholder
    actual_placeholder = get_largest_content_placeholder(slide)
    
    if not actual_placeholder:
        logger.error("No content placeholders for table!")
        return
    
    # Use placeholder dimensions
    left = actual_placeholder.left
    top = actual_placeholder.top
    width = actual_placeholder.width
    height = actual_placeholder.height
    
    # Remove placeholder
    try:
        sp = actual_placeholder.element
        sp.getparent().remove(sp)
    except Exception as e:
        logger.warning(f"Could not remove placeholder: {e}")
    
    try:
        table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
        table = table_shape.table
        
        # Smart column widths
        col_max_lengths = [len(str(headers[col])) for col in range(len(headers))]
        for row_data in rows:
            for col_idx in range(len(headers)):
                if col_idx < len(row_data):
                    col_max_lengths[col_idx] = max(col_max_lengths[col_idx], len(str(row_data[col_idx])))
        
        total_length = sum(col_max_lengths)
        if total_length > 0:
            for col_idx in range(len(headers)):
                proportion = max(col_max_lengths[col_idx] / total_length, 0.1)
                table.columns[col_idx].width = int(width * proportion)
        
        # Headers
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header_text)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    try:
                        run.font.name = THEME_FONT_NAME
                        run.font.size = THEME_FONT_SIZE
                    except Exception:
                        pass
        
        # Rows
        for row_idx, row_data in enumerate(rows, start=1):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < len(headers):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_text)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            try:
                                run.font.name = THEME_FONT_NAME
                                run.font.size = THEME_FONT_SIZE
                            except Exception:
                                pass
        
        logger.info(f"✓ Table added using largest placeholder")
    except Exception as e:
        logger.error(f"Failed to add table: {e}")


def _handle_icons_dynamic(slide, slide_json: dict, layout_capability, analyzer, presentation):
    """
    ENHANCED: Handle pictogram slides using existing SlideDeck AI icon logic.
    Uses icons from GlobalConfig.ICONS_DIR with embeddings-based matching.
    """
    
    if 'bullet_points' not in slide_json:
        logger.warning("No bullet_points for pictogram slide")
        return
    
    items = slide_json['bullet_points']
    
    # Validate all items are icon markers
    for item in items:
        if not isinstance(item, str) or not item.startswith(ICON_BEGINNING_MARKER):
            logger.warning(f"Invalid icon item: {item}")
            return
    
    # Parse icon specifications
    icons_texts = []
    for item in items:
        match = ICONS_REGEX.search(item)
        if match:
            icons_texts.append((match.group(1), match.group(2)))
    
    if not icons_texts:
        logger.error("No valid icon specifications found")
        return
    
    logger.info(f"Creating pictogram slide with {len(icons_texts)} icons")
    
    # Get slide dimensions
    slide_width_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_width
    slide_height_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_height
    
    # Find matching icons using embeddings (existing SlideDeck AI logic)
    fallback_icon_files = ice.find_icons([icon_name for icon_name, _ in icons_texts])
    
    # Layout calculation (horizontal arrangement)
    n_items = len(icons_texts)
    total_width = n_items * ICON_SIZE
    spacing = (Inches(slide_width_inch) - total_width) / (n_items + 1)
    top = INCHES_3  # Below title
    
    # Add each icon with background and text
    for idx, (icon_name, accompanying_text) in enumerate(icons_texts):
        # Try to find icon file
        icon_path = f'{GlobalConfig.ICONS_DIR}/{icon_name}.png'
        
        if not os.path.exists(icon_path):
            # Use fallback from embeddings match
            fallback_name = fallback_icon_files[idx]
            logger.info(f"Using fallback icon '{fallback_name}' for '{icon_name}'")
            icon_path = f'{GlobalConfig.ICONS_DIR}/{fallback_name}.png'
        
        if not os.path.exists(icon_path):
            logger.warning(f"Icon not found: {icon_path}, skipping")
            continue
        
        # Calculate position
        left = spacing + idx * (ICON_SIZE + spacing)
        center = left + ICON_SIZE / 2
        
        # Add background shape (rounded rectangle)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            center - INCHES_0_5,
            top - (ICON_BG_SIZE - ICON_SIZE) / 2,
            INCHES_1, INCHES_1
        )
        shape.fill.solid()
        shape.shadow.inherit = False
        
        # Random color for visual interest
        bg_color = random.choice(ICON_COLORS)
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = bg_color
        
        # Add icon image
        slide.shapes.add_picture(icon_path, left, top, height=ICON_SIZE)
        
        # Add text box below icon
        text_box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left=center - INCHES_2 / 2,
            top=top + ICON_SIZE + INCHES_0_2,
            width=INCHES_2,
            height=INCHES_2
        )
        
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.CENTER
        
        # Add text with formatting
        format_text(text_frame.paragraphs[0], accompanying_text)
        
        # Center text vertically
        text_frame.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
        
        # Transparent background for text box
        text_box.fill.background()
        text_box.line.fill.background()
        text_box.shadow.inherit = False
        
        # Theme color for text
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.theme_color = pptx.enum.dml.MSO_THEME_COLOR.TEXT_2
    
    # Add attribution at bottom
    _add_text_at_bottom(
        slide=slide,
        slide_width_inch=slide_width_inch,
        slide_height_inch=slide_height_inch,
        text='Icons from SlideDeck AI repository',
        hyperlink='https://github.com/barun-saha/slide-deck-ai/tree/main/icons/png128'
    )
    
    # Remove unused placeholders
    remove_unused_placeholders(slide)
    
    logger.info(f"✅ Added {n_items} pictograms to slide")


def _handle_double_column_dynamic(slide, slide_json: dict, layout_capability, analyzer):
    """Handle double column - NO HARDCODED SIZES"""
    if 'bullet_points' not in slide_json:
        return
        
    columns = slide_json['bullet_points']
    text_placeholders = layout_capability.text_placeholders
    
    if len(text_placeholders) < 2:
        logger.error(f"Need 2 placeholders for double column, found {len(text_placeholders)}")
        return
    
    for idx, column_data in enumerate(columns[:2]):
        placeholder_info = text_placeholders[idx]
        
        actual_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == placeholder_info.idx:
                actual_placeholder = shape
                break
        
        if not actual_placeholder or not actual_placeholder.has_text_frame:
            continue
        
        text_frame = actual_placeholder.text_frame
        text_frame.clear()
        
        # Column heading - level 0, bold
        if 'heading' in column_data:
            p = text_frame.paragraphs[0]
            p.level = 0
            run = p.add_run()
            run.text = column_data['heading']
            run.font.bold = True
        
        # Bullet points - level 1
        if 'bullet_points' in column_data:
            flat_items = get_flat_list_of_contents(column_data['bullet_points'], level=0)
            for item_text, item_level in flat_items:
                p = text_frame.add_paragraph()
                p.level = 1
                format_text(p, item_text)
    
    logger.info("✓ Double column populated")

def _handle_bullets_dynamic(slide, slide_json: dict, layout_capability, analyzer):
    """Handle bullets - distribute intelligently across multiple content placeholders.

    Behavior:
    - If layout_capability provides content_placeholders, use those (left-to-right order).
    - If multiple placeholders exist and `bullet_points` is a list of section dicts,
      map each section to its corresponding placeholder (comparison layout).
    - Otherwise split the flat bullet list evenly across available placeholders.
    - Falls back to previous behavior (largest placeholder) if no placeholders found.
    """
    if 'bullet_points' not in slide_json:
        return

    # Prefer layout capability placeholders (they are PlaceholderInfo objects)
    ph_infos = []
    try:
        if layout_capability and getattr(layout_capability, 'content_placeholders', None):
            ph_infos = list(layout_capability.content_placeholders)
    except Exception:
        ph_infos = []

    # If we didn't get placeholders from layout_capability, fall back to shapes left-to-right
    if not ph_infos:
        ph_shapes = get_content_placeholders_left_to_right(slide)
        ph_infos = []
        for sh in ph_shapes:
            try:
                ph_infos.append(type('PH', (), {'idx': sh.placeholder_format.idx, 'left': sh.left, 'area': (sh.width * sh.height)}))
            except Exception:
                continue

    if not ph_infos:
        # Fallback to previous behavior
        actual_placeholder = get_largest_content_placeholder(slide)
        if not actual_placeholder or not actual_placeholder.has_text_frame:
            logger.error("No text placeholders for bullets!")
            return
        text_frame = actual_placeholder.text_frame
        flat_items = get_flat_list_of_contents(slide_json['bullet_points'], level=0)
        add_bulleted_items(text_frame, flat_items)
        logger.info("✓ Bullets added to largest placeholder (fallback)")
        return

    # Sort placeholders left-to-right by 'left' attribute if present
    try:
        ph_infos.sort(key=lambda p: getattr(p, 'left', 0))
    except Exception:
        pass

    # If bullets is a list of dicts with headings (comparison), map per placeholder
    bullets = slide_json['bullet_points']
    if isinstance(bullets, list) and bullets and all(isinstance(b, dict) for b in bullets) and len(bullets) <= len(ph_infos):
        for idx, section_data in enumerate(bullets[:len(ph_infos)]):
            ph = ph_infos[idx]
            placeholder = get_placeholder_safely(slide, getattr(ph, 'idx', ph.idx), placeholder_name='multi-bullets')
            if not placeholder or not placeholder.has_text_frame:
                continue
            tf = placeholder.text_frame
            tf.clear()
            # Heading
            if section_data.get('heading'):
                p = tf.paragraphs[0]
                p.level = 0
                run = p.add_run()
                run.text = section_data.get('heading')
                run.font.bold = True
            flat_items = get_flat_list_of_contents(section_data.get('bullet_points', []), level=0)
            add_bulleted_items(tf, flat_items)
        logger.info("✓ Distributed comparison bullets across placeholders")
        return

    # Otherwise split the flat list across available placeholders
    flat_items = get_flat_list_of_contents(bullets, level=0)
    n_ph = len(ph_infos)
    if n_ph == 0:
        logger.error("No content placeholders available for bullets")
        return

    # Chunk the flat_items into n_ph parts as evenly as possible
    chunks = [[] for _ in range(n_ph)]
    for i, item in enumerate(flat_items):
        chunks[i % n_ph].append(item)

    for idx, ph in enumerate(ph_infos):
        ph_idx = getattr(ph, 'idx', ph.idx)
        placeholder = get_placeholder_safely(slide, ph_idx, placeholder_name=f'col_{idx+1}')
        if not placeholder or not placeholder.has_text_frame:
            continue
        tf = placeholder.text_frame
        tf.clear()
        add_bulleted_items(tf, chunks[idx])

    logger.info(f"✓ Distributed {len(flat_items)} bullets across {n_ph} placeholders")

# --------------------------------------------------------------------------------------------- #

if __name__ == '__main__':
    _JSON_DATA = '''
    {
  "title": "AI Applications: Transforming Industries",
  "slides": [
    {
      "heading": "Introduction to AI Applications",
      "bullet_points": [
        "Artificial Intelligence (AI) is *transforming* various industries",
        "AI applications range from simple decision-making tools to complex systems",
        "AI can be categorized into types: Rule-based, Instance-based, and Model-based"
      ],
      "key_message": "AI is a broad field with diverse applications and categories",
      "img_keywords": "AI, transformation, industries, decision-making, categories"
    },
    {
      "heading": "AI in Everyday Life",
      "bullet_points": [
        "**Virtual assistants** like Siri, Alexa, and Google Assistant",
        "**Recommender systems** in Netflix, Amazon, and Spotify",
        "**Fraud detection** in banking and *credit card* transactions"
      ],
      "key_message": "AI is integrated into our daily lives through various services",
      "img_keywords": "virtual assistants, recommender systems, fraud detection"
    },
    {
      "heading": "AI in Healthcare",
      "bullet_points": [
        "Disease diagnosis and prediction using machine learning algorithms",
        "Personalized medicine and drug discovery",
        "AI-powered robotic surgeries and remote patient monitoring"
      ],
      "key_message": "AI is revolutionizing healthcare with improved diagnostics and patient care",
      "img_keywords": "healthcare, disease diagnosis, personalized medicine, robotic surgeries"
    },
    {
      "heading": "AI in Key Industries",
      "bullet_points": [
        {
          "heading": "Retail",
          "bullet_points": [
            "Inventory management and demand forecasting",
            "Customer segmentation and targeted marketing",
            "AI-driven chatbots for customer service"
          ]
        },
        {
          "heading": "Finance",
          "bullet_points": [
            "Credit scoring and risk assessment",
            "Algorithmic trading and portfolio management",
            "AI for detecting money laundering and cyber fraud"
          ]
        }
      ],
      "key_message": "AI is transforming retail and finance with improved operations and decision-making",
      "img_keywords": "retail, finance, inventory management, credit scoring, algorithmic trading"
    },
    {
      "heading": "AI in Education",
      "bullet_points": [
        "Personalized learning paths and adaptive testing",
        "Intelligent tutoring systems for skill development",
        "AI for predicting student performance and dropout rates"
      ],
      "key_message": "AI is personalizing education and improving student outcomes",
    },
    {
      "heading": "Step-by-Step: AI Development Process",
      "bullet_points": [
        ">> **Step 1:** Define the problem and objectives",
        ">> **Step 2:** Collect and preprocess data",
        ">> **Step 3:** Select and train the AI model",
        ">> **Step 4:** Evaluate and optimize the model",
        ">> **Step 5:** Deploy and monitor the AI system"
      ],
      "key_message": "Developing AI involves a structured process from problem definition to deployment",
      "img_keywords": ""
    },
    {
      "heading": "AI Icons: Key Aspects",
      "bullet_points": [
        "[[brain]] Human-like *intelligence* and decision-making",
        "[[robot]] Automation and physical *tasks*",
        "[[]] Data processing and cloud computing",
        "[[lightbulb]] Insights and *predictions*",
        "[[globe2]] Global connectivity and *impact*"
      ],
      "key_message": "AI encompasses various aspects, from human-like intelligence to global impact",
      "img_keywords": "AI aspects, intelligence, automation, data processing, global impact"
    },
    {
        "heading": "AI vs. ML vs. DL: A Tabular Comparison",
        "table": {
            "headers": ["Feature", "AI", "ML", "DL"],
            "rows": [
                ["Definition", "Creating intelligent agents", "Learning from data", "Deep neural networks"],
                ["Approach", "Rule-based, expert systems", "Algorithms, statistical models", "Deep neural networks"],
                ["Data Requirements", "Varies", "Large datasets", "Massive datasets"],
                ["Complexity", "Varies", "Moderate", "High"],
                ["Computational Cost", "Low to Moderate", "Moderate", "High"],
                ["Examples", "Chess, recommendation systems", "Spam filters, image recognition", "Image recognition, NLP"]
            ]
        },
        "key_message": "This table provides a concise comparison of the key features of AI, ML, and DL.",
        "img_keywords": "AI, ML, DL, comparison, table, features"
    },
    {
      "heading": "Conclusion: Embracing AI's Potential",
      "bullet_points": [
        "AI is transforming industries and improving lives",
        "Ethical considerations are crucial for responsible AI development",
        "Invest in AI education and workforce development",
        "Call to action: Explore AI applications and contribute to shaping its future"
      ],
      "key_message": "AI offers *immense potential*, and we must embrace it responsibly",
      "img_keywords": "AI transformation, ethical considerations, AI education, future of AI"
    }
  ]
}'''

    temp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    path = pathlib.Path(temp.name)

    generate_powerpoint_presentation(
        json5.loads(_JSON_DATA),
        output_file_path=path,
        slides_template='Basic'
    )
    print(f'File path: {path}')

    temp.close()

"""
A set of functions to create a PowerPoint slide deck.
"""
import logging
import os
import pathlib
import random
import re
import tempfile
from typing import Optional

import json5
import pptx
from dotenv import load_dotenv
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.shapes.placeholder import PicturePlaceholder, SlidePlaceholder

from . import icons_embeddings as ice
from . import image_search as ims
from ..global_config import GlobalConfig
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
load_dotenv()

# English Metric Unit (used by PowerPoint) to inches
EMU_TO_INCH_SCALING_FACTOR = 1.0 / 914400
INCHES_3 = pptx.util.Inches(3)
INCHES_2 = pptx.util.Inches(2)
INCHES_1_5 = pptx.util.Inches(1.5)
INCHES_1 = pptx.util.Inches(1)
INCHES_0_8 = pptx.util.Inches(0.8)
INCHES_0_9 = pptx.util.Inches(0.9)
INCHES_0_5 = pptx.util.Inches(0.5)
INCHES_0_4 = pptx.util.Inches(0.4)
INCHES_0_3 = pptx.util.Inches(0.3)
INCHES_0_2 = pptx.util.Inches(0.2)

STEP_BY_STEP_PROCESS_MARKER = '>> '
ICON_BEGINNING_MARKER = '[['
ICON_END_MARKER = ']]'

ICON_SIZE = INCHES_0_8
ICON_BG_SIZE = INCHES_1

IMAGE_DISPLAY_PROBABILITY = 1 / 3.0
FOREGROUND_IMAGE_PROBABILITY = 0.8

SLIDE_NUMBER_REGEX = re.compile(r"^slide[ ]+\d+:", re.IGNORECASE)
ICONS_REGEX = re.compile(r"\[\[(.*?)\]\]\s*(.*)")
BOLD_ITALICS_PATTERN = re.compile(r'(\*\*(.*?)\*\*|\*(.*?)\*)')

ICON_COLORS = [
    pptx.dml.color.RGBColor.from_string('800000'),  # Maroon
    pptx.dml.color.RGBColor.from_string('6A5ACD'),  # SlateBlue
    pptx.dml.color.RGBColor.from_string('556B2F'),  # DarkOliveGreen
    pptx.dml.color.RGBColor.from_string('2F4F4F'),  # DarkSlateGray
    pptx.dml.color.RGBColor.from_string('4682B4'),  # SteelBlue
    pptx.dml.color.RGBColor.from_string('5F9EA0'),  # CadetBlue
]


logger = logging.getLogger(__name__)
logging.getLogger('PIL.PngImagePlugin').setLevel(logging.ERROR)

def remove_unused_placeholders(slide):
    """Remove empty placeholders from a slide (keep title)."""
    shapes_to_remove = []
    
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.idx == 0:  # Keep title
            continue
            
        is_empty = True
        if shape.has_text_frame and shape.text_frame.text.strip():
            is_empty = False
        
        if is_empty:
            shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        try:
            sp = shape.element
            sp.getparent().remove(sp)
        except Exception as e:
            logger.warning(f"Could not remove placeholder: {e}")

def get_content_placeholders_left_to_right(slide):
    """Get content placeholders ordered from left to right (skip title)."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    
    content_placeholders = []
    
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            continue
        try:
            if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                continue
        except:
            pass
        
        content_placeholders.append((shape, shape.left))
    
    if not content_placeholders:
        return []
    
    content_placeholders.sort(key=lambda x: x[1])
    return [ph[0] for ph in content_placeholders]
    
def get_largest_content_placeholder(slide):
    """
    Find the largest non-title placeholder, which is typically the content area.
    
    Args:
        slide: The slide object
    
    Returns:
        The largest content placeholder, or None if not found
    """
    from pptx.enum.shapes import PP_PLACEHOLDER
    
    content_placeholders = []
    
    for shape in slide.placeholders:
        # Skip title placeholders
        if shape.placeholder_format.idx == 0:
            continue
        try:
            if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                continue
        except:
            pass
            
        # Calculate area
        area = shape.width * shape.height
        content_placeholders.append((shape, area))
    
    if not content_placeholders:
        return None
    
    # Return the placeholder with the largest area
    return max(content_placeholders, key=lambda x: x[1])[0]
    
def get_placeholder_safely(slide, idx, placeholder_name=""):
    """
    Safely get a placeholder from a slide by index.
    Falls back to the largest non-title placeholder if index not found.
    
    Args:
        slide: The slide object
        idx: The placeholder index to retrieve
        placeholder_name: Optional name for better error messages
    
    Returns:
        The placeholder object if found, None otherwise
    """
    try:
        return slide.placeholders[idx]
    except KeyError:
        # Get all available placeholder indices
        available_indices = [ph.placeholder_format.idx for ph in slide.placeholders]
        
        logger.warning(
            f"Placeholder {idx} ({placeholder_name}) not found in slide. "
            f"Available indices: {available_indices}"
        )
        
        # Try to find the largest non-title placeholder
        largest_placeholder = get_largest_content_placeholder(slide)
        
        if largest_placeholder:
            logger.info(f"Using largest content placeholder at index {largest_placeholder.placeholder_format.idx}")
            return largest_placeholder
        
        # Last resort: try first available non-title placeholder
        non_title_indices = [idx for idx in available_indices if idx != 0]
        if non_title_indices:
            fallback_idx = non_title_indices[0]
            logger.info(f"Using fallback placeholder at index {fallback_idx}")
            return slide.placeholders[fallback_idx]
        
        logger.error("No placeholders found in this slide at all!")
        return None


def get_placeholder_by_type(slide, placeholder_type):
    """
    Get placeholder by type instead of index (more robust).
    
    Args:
        slide: The slide object
        placeholder_type: Type of placeholder (e.g., PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.BODY)
    
    Returns:
        The first placeholder of the specified type, or None
    """
    for shape in slide.placeholders:
        if shape.placeholder_format.type == placeholder_type:
            return shape
    return None
    
def remove_slide_number_from_heading(header: str) -> str:
    """Remove slide number prefix like 'Slide 1:' from heading."""
    import re
    SLIDE_NUMBER_REGEX = re.compile(r"^slide[ ]+\d+:", re.IGNORECASE)
    if SLIDE_NUMBER_REGEX.match(header):
        idx = header.find(':')
        header = header[idx + 1:].strip()
    return header


def add_bulleted_items(text_frame, flat_items_list: list):
    """
    Add bullet points to a text frame with proper formatting.
    
    Args:
        text_frame: The text frame to add bullets to
        flat_items_list: List of (text, level) tuples
    """
    for idx, an_item in enumerate(flat_items_list):
        if idx == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
            paragraph.level = an_item[1]
        
        STEP_BY_STEP_PROCESS_MARKER = '>> '
        format_text(paragraph, an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER))

def format_text(frame_paragraph, text: str):
    """
    Apply bold (**text**) and italic (*text*) formatting.
    
    Args:
        frame_paragraph: The paragraph to format
        text: Text with markdown-style formatting
    """
    import re
    BOLD_ITALICS_PATTERN = re.compile(r'(\*\*(.*?)\*\*|\*(.*?)\*)')
    
    matches = list(BOLD_ITALICS_PATTERN.finditer(text))
    last_index = 0
    
    for match in matches:
        start, end = match.span()
        if start > last_index:
            run = frame_paragraph.add_run()
            run.text = text[last_index:start]
        
        if match.group(2):  # Bold
            run = frame_paragraph.add_run()
            run.text = match.group(2)
            run.font.bold = True
        elif match.group(3):  # Italics
            run = frame_paragraph.add_run()
            run.text = match.group(3)
            run.font.italic = True
        
        last_index = end
    
    if last_index < len(text):
        run = frame_paragraph.add_run()
        run.text = text[last_index:]

def generate_powerpoint_presentation(
        parsed_data: dict,
        slides_template: str,
        output_file_path: pathlib.Path
) -> list:
    """
    COMPLETE FIXED VERSION
    """
    presentation = pptx.Presentation(GlobalConfig.PPTX_TEMPLATE_FILES[slides_template]['file'])
    
    # Clear existing slides
    slide_id_list = [slide.slide_id for slide in presentation.slides]
    for slide_id in slide_id_list:
        rId = presentation.slides._sldIdLst[0].rId
        presentation.part.drop_rel(rId)
        del presentation.slides._sldIdLst[0]
    
    slide_width_inch, slide_height_inch = _get_slide_width_height_inches(presentation)

    # Title slide
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    presentation_title = parsed_data.get('title', 'Untitled Presentation')
    
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = presentation_title
    
    # Get subtitle placeholder
    subtitle_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0:
            subtitle_shape = shape
            break
    
    if subtitle_shape:
        try:
            subtitle_shape.text = 'by RRD :)'
        except Exception as e:
            logger.warning(f"Error setting subtitle: {e}")
    
    all_headers = [presentation_title]

    # MAIN CONTENT LOOP
    for a_slide in parsed_data['slides']:
        try:
            is_processing_done = False
            current_slide = None
            
            # 1. Try chart
            is_processing_done = _handle_chart(
                presentation=presentation,
                slide_json=a_slide,
                slide_width_inch=slide_width_inch,
                slide_height_inch=slide_height_inch
            )
            if is_processing_done:
                current_slide = presentation.slides[-1]

            # 2. Try icons
            if not is_processing_done:
                is_processing_done = _handle_icons_ideas(
                    presentation=presentation,
                    slide_json=a_slide,
                    slide_width_inch=slide_width_inch,
                    slide_height_inch=slide_height_inch
                )
                if is_processing_done:
                    current_slide = presentation.slides[-1]

            # 3. Try table
            if not is_processing_done:
                is_processing_done = _handle_table(
                    presentation=presentation,
                    slide_json=a_slide,
                    slide_width_inch=slide_width_inch,
                    slide_height_inch=slide_height_inch
                )
                if is_processing_done:
                    current_slide = presentation.slides[-1]

            # 4. Try double column
            if not is_processing_done:
                is_processing_done = _handle_double_col_layout(
                    presentation=presentation,
                    slide_json=a_slide,
                    slide_width_inch=slide_width_inch,
                    slide_height_inch=slide_height_inch
                )
                if is_processing_done:
                    current_slide = presentation.slides[-1]

            # 5. Try step by step
            if not is_processing_done:
                is_processing_done = _handle_step_by_step_process(
                    presentation=presentation,
                    slide_json=a_slide,
                    slide_width_inch=slide_width_inch,
                    slide_height_inch=slide_height_inch
                )
                if is_processing_done:
                    current_slide = presentation.slides[-1]

            # 6. Default display
            if not is_processing_done:
                current_slide = _handle_default_display(
                    presentation=presentation,
                    slide_json=a_slide,
                    slide_width_inch=slide_width_inch,
                    slide_height_inch=slide_height_inch
                )
            
            # *** ADD KEY MESSAGE TO EVERY SLIDE ***
            if current_slide and 'key_message' in a_slide and a_slide['key_message']:
                _handle_key_message(
                    the_slide=current_slide,
                    slide_json=a_slide,
                    slide_width_inch=slide_width_inch,
                    slide_height_inch=slide_height_inch
                )

        except Exception:
            logger.error(
                'An error occurred while processing a slide',
                exc_info=True
            )
            continue

    # Thank you slide
    last_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(last_slide_layout)
    
    title_shape = slide.shapes.title
    if title_shape:
        try:
            title_shape.text = 'Thank you!'
        except Exception as e:
            logger.warning(f"Error setting thank you title: {e}")

    presentation.save(output_file_path)
    return all_headers


def get_flat_list_of_contents(items: list, level: int) -> list:
    """
    Flatten hierarchical bullet points into (text, level) tuples.
    
    Args:
        items: List of strings or nested lists
        level: Current hierarchy level
    
    Returns:
        List of (text, level) tuples
    """
    flat_list = []
    for item in items:
        if isinstance(item, str):
            flat_list.append((item, level))
        elif isinstance(item, list):
            flat_list = flat_list + get_flat_list_of_contents(item, level + 1)
    return flat_list


def get_slide_placeholders(
        slide: pptx.slide.Slide,
        layout_number: int,
        is_debug: bool = False
) -> list[tuple[int, str]]:
    """
    Return the index and name (lower case) of all placeholders present in a
    slide, except the title placeholder.

    A placeholder in a slide is a place to add content. Each placeholder has a
    name and an index. This index is not a list index; it is a key used to look up
    a dict and may be non-contiguous. The title placeholder always has index 0.
    User-added placeholders get indices starting from 10.

    With user-edited or added placeholders, indices may be difficult to track. This
    function returns the placeholders' names as well, which may help distinguish
    between placeholders.

    Args:
        slide: The slide.
        layout_number: The layout number used by the slide.
        is_debug: Whether to print debugging statements.

    Returns:
        list[tuple[int, str]]: A list of (index, name) tuples for placeholders
        present in the slide, excluding the title placeholder.
    """

    if is_debug:
        print(
            f'Slide layout #{layout_number}:'
            f' # of placeholders: {len(slide.shapes.placeholders)} (including the title)'
        )

    placeholders = [
        (shape.placeholder_format.idx, shape.name.lower()) for shape in slide.shapes.placeholders
    ]
    placeholders.pop(0)  # Remove the title placeholder

    if is_debug:
        print(placeholders)

    return placeholders

def _handle_chart(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """
    FIXED: Charts now use MAXIMUM available space, not just placeholder size
    FIXED: Proper handling of cx/cy for logging
    """
    if 'chart' not in slide_json or not slide_json['chart']:
        return False
    
    chart_data_json = slide_json['chart']
    chart_type = chart_data_json.get('type', 'column').lower()
    categories = chart_data_json.get('categories', [])
    series_data = chart_data_json.get('series', [])
    
    if not categories or not series_data:
        logger.warning("Chart data incomplete - missing categories or series")
        return False
    
    # Use a layout with title and content
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = remove_slide_number_from_heading(slide_json['heading'])
    
    # Get LEFTMOST placeholder
    placeholders = get_content_placeholders_left_to_right(slide)
    
    # *** CRITICAL FIX: Use MAXIMUM space, not placeholder size ***
    if placeholders and len(placeholders) > 0:
        content_placeholder = placeholders[0]
        
        # Remove placeholder first
        try:
            sp = content_placeholder.element
            sp.getparent().remove(sp)
        except Exception as e:
            logger.warning(f"Could not remove placeholder: {e}")
        
        # Use generous margins for chart
        margin = Inches(0.5)
        title_height = Inches(1.5) if title_shape else Inches(1.0)
        
        x = margin
        y = title_height
        cx = Inches(slide_width_inch) - (2 * margin)  # Full width minus margins
        cy = Inches(slide_height_inch) - title_height - margin - Inches(0.5)  # Max height
    else:
        # Fallback: use maximum space
        x = Inches(0.5)
        y = Inches(1.5)
        cx = Inches(slide_width_inch - 1)
        cy = Inches(slide_height_inch - 2.5)
    
    # Determine chart type
    chart_type_map = {
        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'line': XL_CHART_TYPE.LINE,
        'pie': XL_CHART_TYPE.PIE,
    }
    
    xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # Prepare chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    
    if chart_type == 'pie':
        if isinstance(series_data[0], dict):
            series_name = series_data[0].get('name', 'Values')
            series_values = series_data[0].get('values', [])
        else:
            series_name = 'Values'
            series_values = series_data
        chart_data.add_series(series_name, series_values)
    else:
        for series_item in series_data:
            if isinstance(series_item, dict):
                series_name = series_item.get('name', 'Series')
                series_values = series_item.get('values', [])
            else:
                series_name = 'Series'
                series_values = series_item
            chart_data.add_series(series_name, series_values)
    
    # Add chart with MAXIMUM dimensions
    try:
        chart = slide.shapes.add_chart(
            xl_chart_type,
            x, y, cx, cy,
            chart_data
        ).chart
        
        chart.has_legend = True
        if chart_type != 'pie':
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        
        if 'title' in chart_data_json:
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_data_json['title']
        
        # FIXED: Safe logging without assuming .inches attribute
        try:
            width_inches = cx.inches if hasattr(cx, 'inches') else cx / 914400
            height_inches = cy.inches if hasattr(cy, 'inches') else cy / 914400
            logger.info(f"Successfully added {chart_type} chart (size: {width_inches:.1f}\" x {height_inches:.1f}\")")
        except Exception:
            logger.info(f"Successfully added {chart_type} chart")
        
    except Exception as ex:
        logger.error(f"Error adding chart to slide: {str(ex)}", exc_info=True)
        return False
    
    remove_unused_placeholders(slide)
    return True
    
def _handle_default_display(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
):
    """
    FIXED: Added proper validation for 'bullet_points' key existence.
    
    Handle default bullet point slides with optional image display.
    
    Args:
        presentation: The PowerPoint presentation object
        slide_json: Dictionary containing slide data
        slide_width_inch: Width of slide in inches
        slide_height_inch: Height of slide in inches
    
    Returns:
        The created slide object
    """
    import random
    
    status = False
    slide = None

    # Try to add images if keywords are present
    if 'img_keywords' in slide_json:
        if random.random() < IMAGE_DISPLAY_PROBABILITY:
            if random.random() < FOREGROUND_IMAGE_PROBABILITY:
                status, slide = _handle_display_image__in_foreground(
                    presentation, slide_json, slide_width_inch, slide_height_inch
                )
            else:
                status, slide = _handle_display_image__in_background(
                    presentation, slide_json, slide_width_inch, slide_height_inch
                )

    if status and slide:
        return slide

    # Create bullet slide
    bullet_slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    if title_shape:
        title_shape.text = remove_slide_number_from_heading(slide_json['heading'])
    
    # Get LEFTMOST placeholder
    placeholders = get_content_placeholders_left_to_right(slide)
    
    body_shape = None
    if placeholders and len(placeholders) > 0:
        body_shape = placeholders[0]
    
    # FIX: Check if 'bullet_points' exists before accessing
    if body_shape and body_shape.has_text_frame and 'bullet_points' in slide_json:
        text_frame = body_shape.text_frame
        flat_items_list = get_flat_list_of_contents(slide_json['bullet_points'], level=0)
        add_bulleted_items(text_frame, flat_items_list)
    elif body_shape and body_shape.has_text_frame:
        # If no bullet points, add a placeholder message
        text_frame = body_shape.text_frame
        text_frame.text = "(Content not available)"
    
    remove_unused_placeholders(slide)
    return slide
    
def _handle_display_image__in_foreground(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> tuple[bool, pptx.slide.Slide]:
    """Returns (status, slide) tuple"""
    
    img_keywords = slide_json.get('img_keywords', '').strip()
    
    slide_layout = presentation.slide_layouts[8]
    slide = presentation.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    if title_placeholder:
        title_placeholder.text = remove_slide_number_from_heading(slide_json['heading'])
    
    placeholders = get_content_placeholders_left_to_right(slide)
    pic_col = placeholders[0] if len(placeholders) > 0 else None
    text_col = placeholders[1] if len(placeholders) > 1 else None
    
    if text_col and text_col.has_text_frame and 'bullet_points' in slide_json:
        flat_items_list = get_flat_list_of_contents(slide_json['bullet_points'], level=0)
        add_bulleted_items(text_col.text_frame, flat_items_list)
    
    if img_keywords and pic_col:
        try:
            photo_url, page_url = ims.get_photo_url_from_api_response(
                ims.search_pexels(query=img_keywords, size='medium')
            )
            if photo_url:
                pic_col.insert_picture(ims.get_image_from_url(photo_url))
                _add_text_at_bottom(
                    slide=slide,
                    slide_width_inch=slide_width_inch,
                    slide_height_inch=slide_height_inch,
                    text='Photo provided by Pexels',
                    hyperlink=page_url
                )
        except Exception as ex:
            logger.error(f'Error adding image: {ex}')
    
    remove_unused_placeholders(slide)
    return True, slide


def _handle_display_image__in_background(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> tuple[bool, pptx.slide.Slide]:
    """Returns (status, slide) tuple"""
    
    img_keywords = slide_json.get('img_keywords', '').strip()
    
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title_shape = slide.shapes.title
    
    placeholders = get_content_placeholders_left_to_right(slide)
    body_shape = placeholders[0] if placeholders else None
    
    if title_shape:
        title_shape.text = remove_slide_number_from_heading(slide_json['heading'])
    
    if body_shape and 'bullet_points' in slide_json:
        flat_items_list = get_flat_list_of_contents(slide_json['bullet_points'], level=0)
        add_bulleted_items(body_shape.text_frame, flat_items_list)
    
    if img_keywords:
        try:
            photo_url, page_url = ims.get_photo_url_from_api_response(
                ims.search_pexels(query=img_keywords, size='large')
            )
            if photo_url:
                picture = slide.shapes.add_picture(
                    image_file=ims.get_image_from_url(photo_url),
                    left=0, top=0,
                    width=pptx.util.Inches(slide_width_inch),
                )
                
                # Apply transparency
                try:
                    blip_elements = picture._element.xpath('.//a:blip')
                    for blip in blip_elements:
                        alpha_mod = blip.makeelement(
                            '{http://schemas.openxmlformats.org/drawingml/2006/main}alphaModFix'
                        )
                        alpha_mod.set('amt', '50000')
                        existing = blip.find(
                            '{http://schemas.openxmlformats.org/drawingml/2006/main}alphaModFix'
                        )
                        if existing is not None:
                            blip.remove(existing)
                        blip.append(alpha_mod)
                except Exception as ex:
                    logger.error(f'Failed to apply transparency: {ex}')
                
                _add_text_at_bottom(
                    slide=slide,
                    slide_width_inch=slide_width_inch,
                    slide_height_inch=slide_height_inch,
                    text='Photo provided by Pexels',
                    hyperlink=page_url
                )
                
                # Move to background
                try:
                    slide.shapes._spTree.remove(picture._element)
                    slide.shapes._spTree.insert(2, picture._element)
                except Exception as ex:
                    logger.error(f'Failed to move to background: {ex}')
        except Exception as ex:
            logger.error(f'Error adding background image: {ex}')
    
    remove_unused_placeholders(slide)
    return True, slide


def _handle_icons_ideas(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
):
    """
    Add a slide with some icons and text.
    If no suitable icons are found, the step numbers are shown.

    Args:
        presentation: The presentation object.
        slide_json: The content of the slide as JSON data.
        slide_width_inch: The width of the slide in inches.
        slide_height_inch: The height of the slide in inches.

    Returns:
        True if the slide has been processed.
    """

    if 'bullet_points' in slide_json and slide_json['bullet_points']:
        items = slide_json['bullet_points']

        # Ensure that it is a single list of strings without any sub-list
        for step in items:
            if not isinstance(step, str) or not step.startswith(ICON_BEGINNING_MARKER):
                return False

        slide_layout = presentation.slide_layouts[5]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        if title is not None:
            title.text = remove_slide_number_from_heading(slide_json['heading'])
        else:
            logger.warning("No title placeholder in layout 5 (icons)")

        n_items = len(items)
        text_box_size = INCHES_2

        total_width = n_items * ICON_SIZE
        spacing = (pptx.util.Inches(slide_width_inch) - total_width) / (n_items + 1)
        top = INCHES_3

        icons_texts = [
            (match.group(1), match.group(2)) for match in [
                ICONS_REGEX.search(item) for item in items
            ]
        ]
        fallback_icon_files = ice.find_icons([item[0] for item in icons_texts])

        for idx, item in enumerate(icons_texts):
            icon, accompanying_text = item
            icon_path = f'{GlobalConfig.ICONS_DIR}/{icon}.png'

            if not os.path.exists(icon_path):
                logger.warning(
                    'Icon not found: %s...using fallback icon: %s',
                    icon, fallback_icon_files[idx]
                )
                icon_path = f'{GlobalConfig.ICONS_DIR}/{fallback_icon_files[idx]}.png'

            left = spacing + idx * (ICON_SIZE + spacing)
            center = left + ICON_SIZE / 2

            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                center - INCHES_0_5,
                top - (ICON_BG_SIZE - ICON_SIZE) / 2,
                INCHES_1, INCHES_1
            )
            shape.fill.solid()
            shape.shadow.inherit = False
            shape.fill.fore_color.rgb = shape.line.color.rgb = random.choice(ICON_COLORS)
            
            slide.shapes.add_picture(icon_path, left, top, height=ICON_SIZE)

            text_box = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                left=center - text_box_size / 2,
                top=top + ICON_SIZE + INCHES_0_2,
                width=text_box_size,
                height=text_box_size
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.CENTER
            format_text(text_frame.paragraphs[0], accompanying_text)

            text_frame.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
            text_box.fill.background()
            text_box.line.fill.background()
            text_box.shadow.inherit = False

            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.theme_color = pptx.enum.dml.MSO_THEME_COLOR.TEXT_2

        _add_text_at_bottom(
            slide=slide,
            slide_width_inch=slide_width_inch,
            slide_height_inch=slide_height_inch,
            text='More icons available in the SlideDeck AI repository',
            hyperlink='https://github.com/barun-saha/slide-deck-ai/tree/main/icons/png128'
        )

        remove_unused_placeholders(slide)
        return True
    
    return False


def _add_text_at_bottom(
        slide: pptx.slide.Slide,
        slide_width_inch: float,
        slide_height_inch: float,
        text: str,
        hyperlink: Optional[str] = None,
        target_height: Optional[float] = 0.5
):
    """
    Add arbitrary text to a textbox positioned near the lower-left side of a slide.

    Args:
        slide: The slide.
        slide_width_inch: The width of the slide in inches.
        slide_height_inch: The height of the slide in inches.
        text: The text to be added.
        hyperlink: Optional; the hyperlink to be added to the text.
        target_height: Optional[float]; the target height of the box in inches.
    """

    footer = slide.shapes.add_textbox(
        left=INCHES_1,
        top=pptx.util.Inches(slide_height_inch - target_height),
        width=pptx.util.Inches(slide_width_inch),
        height=pptx.util.Inches(target_height)
    )

    paragraph = footer.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = pptx.util.Pt(10)
    run.font.underline = False

    if hyperlink:
        run.hyperlink.address = hyperlink
        
def _handle_double_col_layout(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """
    FIXED: Proper text hierarchy with headings as bold level 0 and bullets as level 1.
    FIXED: Column headings now appear INSIDE content area, not overlaying slide title.
    
    Create a two-column layout slide with proper formatting:
    - Slide title: Main heading (e.g., "Strengths vs. Challenges")
    - Column headings: Bold, 16pt, level 0 (no bullet) INSIDE content placeholders
    - Bullet points: Normal, level 1 (indented with bullets)
    
    Args:
        presentation: The PowerPoint presentation object
        slide_json: Dictionary containing slide data with double-column structure
        slide_width_inch: Width of slide in inches
        slide_height_inch: Height of slide in inches
    
    Returns:
        bool: True if successfully created, False otherwise
    """
    if 'bullet_points' not in slide_json or not slide_json['bullet_points']:
        return False
    
    double_col_content = slide_json['bullet_points']
    
    # Validate structure: must be exactly 2 dict items
    if not (len(double_col_content) == 2 
            and isinstance(double_col_content[0], dict) 
            and isinstance(double_col_content[1], dict)):
        return False
    
    # Use layout 4 (typically a two-column layout)
    slide_layout = presentation.slide_layouts[4]
    slide = presentation.slides.add_slide(slide_layout)
    
    # Set main slide title (e.g., "Strengths vs. Challenges for Apple")
    title_placeholder = slide.shapes.title
    if title_placeholder:
        heading_text = remove_slide_number_from_heading(slide_json['heading'])
        
        # Remove "Double Column:" prefix if present
        if heading_text.lower().startswith('double column:'):
            heading_text = heading_text[14:].strip()
        
        title_placeholder.text = heading_text
    
    # Get content placeholders ordered left-to-right
    placeholders = get_content_placeholders_left_to_right(slide)
    
    left_col = placeholders[0] if len(placeholders) > 0 else None
    right_col = placeholders[1] if len(placeholders) > 1 else None
    
    def populate_column(placeholder, content_dict):
        """
        Populate a column with heading (level 0, bold) and bullets (level 1).
        
        Args:
            placeholder: The placeholder shape to populate
            content_dict: Dictionary with 'heading' and 'bullet_points' keys
        """
        if not placeholder or not placeholder.has_text_frame:
            return
        
        text_frame = placeholder.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        
        # Add column heading (level 0, bold, 16pt, NO BULLET)
        if 'heading' in content_dict:
            p = text_frame.paragraphs[0]
            p.level = 0  # No indentation
            p.space_after = Pt(12)  # Add spacing after heading
            run = p.add_run()
            run.text = content_dict['heading']
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.theme_color = pptx.enum.dml.MSO_THEME_COLOR.TEXT_1
        
        # Add bullet points (level 1, indented WITH bullets)
        if 'bullet_points' in content_dict:
            flat_items = get_flat_list_of_contents(
                content_dict['bullet_points'], level=0
            )
            for item_text, item_level in flat_items:
                p = text_frame.add_paragraph()
                p.level = 1  # Indented under heading with bullet
                format_text(p, item_text)
                # Ensure proper font size for bullets
                for run in p.runs:
                    if not run.font.size:
                        run.font.size = Pt(14)
    
    # Populate both columns
    populate_column(left_col, double_col_content[0])
    populate_column(right_col, double_col_content[1])
    
    remove_unused_placeholders(slide)
    return True

def _handle_step_by_step_process(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """Add shapes to display a step-by-step process in the slide, if available.

    Args:
        presentation (pptx.Presentation): The presentation object.
        slide_json (dict): The content of the slide as JSON data.
        slide_width_inch (float): The width of the slide in inches.
        slide_height_inch (float): The height of the slide in inches.

    Returns:
        bool: True if this slide has a step-by-step process depiction added; False otherwise.
    """

    if 'bullet_points' in slide_json and slide_json['bullet_points']:
        steps = slide_json['bullet_points']

        no_marker_count = 0.0
        n_steps = len(steps)

        for step in steps:
            if not isinstance(step, str):
                return False

            if not step.startswith(STEP_BY_STEP_PROCESS_MARKER):
                no_marker_count += 1

        slide_header = slide_json['heading'].lower()
        if (no_marker_count / n_steps > 0.25) and not (
                ('step-by-step' in slide_header) or ('step by step' in slide_header)
        ):
            return False

        if n_steps < 3 or n_steps > 6:
            return False

        bullet_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title = shapes.title
        if title is not None:
            title.text = remove_slide_number_from_heading(slide_json['heading'])
        else:
            logger.warning("No title placeholder in step-by-step layout")

        if 3 <= n_steps <= 4:
            # Horizontal display
            height = INCHES_1_5
            width = pptx.util.Inches(slide_width_inch / n_steps - 0.01)
            top = pptx.util.Inches(slide_height_inch / 2)
            left = pptx.util.Inches((slide_width_inch - width.inches * n_steps) / 2 + 0.05)

            for step in steps:
                shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
                text_frame = shape.text_frame
                text_frame.clear()
                paragraph = text_frame.paragraphs[0]
                paragraph.alignment = pptx.enum.text.PP_ALIGN.LEFT
                format_text(paragraph, step.removeprefix(STEP_BY_STEP_PROCESS_MARKER))
                left += width - INCHES_0_4
        elif 4 < n_steps <= 6:
            # Vertical display
            height = pptx.util.Inches(0.65)
            top = pptx.util.Inches(slide_height_inch / 4)
            left = INCHES_1

            width = pptx.util.Inches(slide_width_inch * 2 / 3)
            lengths = [len(step) for step in steps]
            font_size_20pt = pptx.util.Pt(20)
            widths = sorted(
                [
                    min(
                        pptx.util.Inches(font_size_20pt.inches * a_len),
                        width
                    ) for a_len in lengths
                ]
            )
            width = widths[len(widths) // 2]

            for step in steps:
                shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON, left, top, width, height)
                text_frame = shape.text_frame
                text_frame.clear()
                paragraph = text_frame.paragraphs[0]
                paragraph.alignment = pptx.enum.text.PP_ALIGN.LEFT
                format_text(paragraph, step.removeprefix(STEP_BY_STEP_PROCESS_MARKER))
                top += height + INCHES_0_3
                left += INCHES_0_5

        remove_unused_placeholders(slide)
        return True

    return False

def _handle_table(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """Smart column widths based on content"""
    if 'table' not in slide_json or not slide_json['table']:
        return False
    
    headers = slide_json['table'].get('headers', [])
    rows = slide_json['table'].get('rows', [])
    
    if not headers or not rows:
        return False
    
    bullet_slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title = shapes.title
    if title:
        title.text = remove_slide_number_from_heading(slide_json['heading'])
    
    placeholders = get_content_placeholders_left_to_right(slide)
    
    if placeholders and len(placeholders) > 0:
        content_placeholder = placeholders[0]
        left = content_placeholder.left
        top = content_placeholder.top
        width = content_placeholder.width
        height = content_placeholder.height
        
        try:
            sp = content_placeholder.element
            sp.getparent().remove(sp)
        except Exception as e:
            logger.warning(f"Could not remove placeholder: {e}")
    else:
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(slide_width_inch - 1)
        height = Inches(slide_height_inch - 3)
    
    try:
        table_shape = slide.shapes.add_table(
            len(rows) + 1, len(headers),
            left, top, width, height
        )
        table = table_shape.table
        
        # Smart column widths
        col_max_lengths = []
        for col_idx in range(len(headers)):
            max_len = len(str(headers[col_idx]))
            for row_data in rows:
                if col_idx < len(row_data):
                    max_len = max(max_len, len(str(row_data[col_idx])))
            col_max_lengths.append(max_len)
        
        total_length = sum(col_max_lengths)
        if total_length > 0:
            for col_idx in range(len(headers)):
                proportion = col_max_lengths[col_idx] / total_length
                proportion = max(proportion, 0.1)
                table.columns[col_idx].width = int(width * proportion)
        
        # Headers
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header_text)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(12)
        
        # Rows
        for row_idx, row_data in enumerate(rows, start=1):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < len(headers):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_text)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(11)
        
    except Exception as ex:
        logger.error(f"Error creating table: {ex}", exc_info=True)
        return False
    
    remove_unused_placeholders(slide)
    return True

def _handle_key_message(
        the_slide: pptx.slide.Slide,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
):
    """
    FIXED: Adaptive font color based on background brightness.
    
    Add a key message box at the bottom center of the slide.
    Text color automatically adapts:
    - Light backgrounds → Black text
    - Dark backgrounds → White text
    
    Args:
        the_slide: The slide object to add the key message to
        slide_json: Dictionary containing 'key_message' key
        slide_width_inch: Width of slide in inches
        slide_height_inch: Height of slide in inches
    """
    if 'key_message' not in slide_json or not slide_json['key_message']:
        return
    
    key_message_text = slide_json['key_message']
    
    # Calculate position and size
    height = Inches(1.6)
    width = Inches(slide_width_inch / 2.3)
    top = Inches(slide_height_inch - height.inches - 0.1)
    left = Inches((slide_width_inch - width.inches) / 2)
    
    # Create the rounded rectangle shape
    shape = the_slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left=left,
        top=top,
        width=width,
        height=height
    )
    
    # Configure text frame
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    # Add the key message text
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = pptx.enum.text.PP_ALIGN.CENTER
    format_text(paragraph, key_message_text)
    
    # Center the text vertically
    text_frame.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
    
    # Style the shape with light gray background
    shape.fill.solid()
    bg_color = pptx.dml.color.RGBColor(240, 240, 240)  # Light gray
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = pptx.dml.color.RGBColor(200, 200, 200)
    shape.line.width = Pt(1)
    
    # FIX: Determine text color based on background brightness
    # Using relative luminance formula (ITU-R BT.709)
    r, g, b = 240, 240, 240  # Current background color
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    
    # Use black text for light backgrounds, white for dark
    if luminance > 0.5:
        text_color = pptx.dml.color.RGBColor(0, 0, 0)  # Black
    else:
        text_color = pptx.dml.color.RGBColor(255, 255, 255)  # White
    
    # Apply text formatting
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.font.size:
                run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = text_color


def _get_slide_width_height_inches(presentation: pptx.Presentation) -> tuple[float, float]:
    """
    Get the dimensions of a slide in inches.

    Args:
        presentation: The presentation object.

    Returns:
        The width and the height.
    """

    slide_width_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_width
    slide_height_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_height

    return slide_width_inch, slide_height_inch


if __name__ == '__main__':
    _JSON_DATA = '''
    {
  "title": "AI Applications: Transforming Industries",
  "slides": [
    {
      "heading": "Introduction to AI Applications",
      "bullet_points": [
        "Artificial Intelligence (AI) is *transforming* various industries",
        "AI applications range from simple decision-making tools to complex systems",
        "AI can be categorized into types: Rule-based, Instance-based, and Model-based"
      ],
      "key_message": "AI is a broad field with diverse applications and categories",
      "img_keywords": "AI, transformation, industries, decision-making, categories"
    },
    {
      "heading": "AI in Everyday Life",
      "bullet_points": [
        "**Virtual assistants** like Siri, Alexa, and Google Assistant",
        "**Recommender systems** in Netflix, Amazon, and Spotify",
        "**Fraud detection** in banking and *credit card* transactions"
      ],
      "key_message": "AI is integrated into our daily lives through various services",
      "img_keywords": "virtual assistants, recommender systems, fraud detection"
    },
    {
      "heading": "AI in Healthcare",
      "bullet_points": [
        "Disease diagnosis and prediction using machine learning algorithms",
        "Personalized medicine and drug discovery",
        "AI-powered robotic surgeries and remote patient monitoring"
      ],
      "key_message": "AI is revolutionizing healthcare with improved diagnostics and patient care",
      "img_keywords": "healthcare, disease diagnosis, personalized medicine, robotic surgeries"
    },
    {
      "heading": "AI in Key Industries",
      "bullet_points": [
        {
          "heading": "Retail",
          "bullet_points": [
            "Inventory management and demand forecasting",
            "Customer segmentation and targeted marketing",
            "AI-driven chatbots for customer service"
          ]
        },
        {
          "heading": "Finance",
          "bullet_points": [
            "Credit scoring and risk assessment",
            "Algorithmic trading and portfolio management",
            "AI for detecting money laundering and cyber fraud"
          ]
        }
      ],
      "key_message": "AI is transforming retail and finance with improved operations and decision-making",
      "img_keywords": "retail, finance, inventory management, credit scoring, algorithmic trading"
    },
    {
      "heading": "AI in Education",
      "bullet_points": [
        "Personalized learning paths and adaptive testing",
        "Intelligent tutoring systems for skill development",
        "AI for predicting student performance and dropout rates"
      ],
      "key_message": "AI is personalizing education and improving student outcomes",
    },
    {
      "heading": "Step-by-Step: AI Development Process",
      "bullet_points": [
        ">> **Step 1:** Define the problem and objectives",
        ">> **Step 2:** Collect and preprocess data",
        ">> **Step 3:** Select and train the AI model",
        ">> **Step 4:** Evaluate and optimize the model",
        ">> **Step 5:** Deploy and monitor the AI system"
      ],
      "key_message": "Developing AI involves a structured process from problem definition to deployment",
      "img_keywords": ""
    },
    {
      "heading": "AI Icons: Key Aspects",
      "bullet_points": [
        "[[brain]] Human-like *intelligence* and decision-making",
        "[[robot]] Automation and physical *tasks*",
        "[[]] Data processing and cloud computing",
        "[[lightbulb]] Insights and *predictions*",
        "[[globe2]] Global connectivity and *impact*"
      ],
      "key_message": "AI encompasses various aspects, from human-like intelligence to global impact",
      "img_keywords": "AI aspects, intelligence, automation, data processing, global impact"
    },
    {
        "heading": "AI vs. ML vs. DL: A Tabular Comparison",
        "table": {
            "headers": ["Feature", "AI", "ML", "DL"],
            "rows": [
                ["Definition", "Creating intelligent agents", "Learning from data", "Deep neural networks"],
                ["Approach", "Rule-based, expert systems", "Algorithms, statistical models", "Deep neural networks"],
                ["Data Requirements", "Varies", "Large datasets", "Massive datasets"],
                ["Complexity", "Varies", "Moderate", "High"],
                ["Computational Cost", "Low to Moderate", "Moderate", "High"],
                ["Examples", "Chess, recommendation systems", "Spam filters, image recognition", "Image recognition, NLP"]
            ]
        },
        "key_message": "This table provides a concise comparison of the key features of AI, ML, and DL.",
        "img_keywords": "AI, ML, DL, comparison, table, features"
    },
    {
      "heading": "Conclusion: Embracing AI's Potential",
      "bullet_points": [
        "AI is transforming industries and improving lives",
        "Ethical considerations are crucial for responsible AI development",
        "Invest in AI education and workforce development",
        "Call to action: Explore AI applications and contribute to shaping its future"
      ],
      "key_message": "AI offers *immense potential*, and we must embrace it responsibly",
      "img_keywords": "AI transformation, ethical considerations, AI education, future of AI"
    }
  ]
}'''

    temp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    path = pathlib.Path(temp.name)

    generate_powerpoint_presentation(
        json5.loads(_JSON_DATA),
        output_file_path=path,
        slides_template='Basic'
    )
    print(f'File path: {path}')

    temp.close()

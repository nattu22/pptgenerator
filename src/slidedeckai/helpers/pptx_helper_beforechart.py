"""
A set of functions to create a PowerPoint slide deck.
"""
import logging
import os
import pathlib
import random
import re
import tempfile
from typing import Optional

import json5
import pptx
from dotenv import load_dotenv
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.shapes.placeholder import PicturePlaceholder, SlidePlaceholder

from . import icons_embeddings as ice
from . import image_search as ims
from ..global_config import GlobalConfig
from pptx.util import Inches, Pt

load_dotenv()

# English Metric Unit (used by PowerPoint) to inches
EMU_TO_INCH_SCALING_FACTOR = 1.0 / 914400
INCHES_3 = pptx.util.Inches(3)
INCHES_2 = pptx.util.Inches(2)
INCHES_1_5 = pptx.util.Inches(1.5)
INCHES_1 = pptx.util.Inches(1)
INCHES_0_8 = pptx.util.Inches(0.8)
INCHES_0_9 = pptx.util.Inches(0.9)
INCHES_0_5 = pptx.util.Inches(0.5)
INCHES_0_4 = pptx.util.Inches(0.4)
INCHES_0_3 = pptx.util.Inches(0.3)
INCHES_0_2 = pptx.util.Inches(0.2)

STEP_BY_STEP_PROCESS_MARKER = '>> '
ICON_BEGINNING_MARKER = '[['
ICON_END_MARKER = ']]'

ICON_SIZE = INCHES_0_8
ICON_BG_SIZE = INCHES_1

IMAGE_DISPLAY_PROBABILITY = 1 / 3.0
FOREGROUND_IMAGE_PROBABILITY = 0.8

SLIDE_NUMBER_REGEX = re.compile(r"^slide[ ]+\d+:", re.IGNORECASE)
ICONS_REGEX = re.compile(r"\[\[(.*?)\]\]\s*(.*)")
BOLD_ITALICS_PATTERN = re.compile(r'(\*\*(.*?)\*\*|\*(.*?)\*)')

ICON_COLORS = [
    pptx.dml.color.RGBColor.from_string('800000'),  # Maroon
    pptx.dml.color.RGBColor.from_string('6A5ACD'),  # SlateBlue
    pptx.dml.color.RGBColor.from_string('556B2F'),  # DarkOliveGreen
    pptx.dml.color.RGBColor.from_string('2F4F4F'),  # DarkSlateGray
    pptx.dml.color.RGBColor.from_string('4682B4'),  # SteelBlue
    pptx.dml.color.RGBColor.from_string('5F9EA0'),  # CadetBlue
]


logger = logging.getLogger(__name__)
logging.getLogger('PIL.PngImagePlugin').setLevel(logging.ERROR)

def remove_unused_placeholders(slide: pptx.slide.Slide):
    """Remove empty placeholders from a slide."""
    shapes_to_remove = []
    
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.idx == 0:  # Keep title
            continue
            
        # Check if empty
        is_empty = True
        if shape.has_text_frame and shape.text_frame.text.strip():
            is_empty = False
        
        if is_empty:
            shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        try:
            sp = shape.element
            sp.getparent().remove(sp)
        except Exception as e:
            logger.warning(f"Could not remove placeholder: {e}")

def get_largest_content_placeholder(slide):
    """
    Find the largest non-title placeholder, which is typically the content area.
    
    Args:
        slide: The slide object
    
    Returns:
        The largest content placeholder, or None if not found
    """
    from pptx.enum.shapes import PP_PLACEHOLDER
    
    content_placeholders = []
    
    for shape in slide.placeholders:
        # Skip title placeholders
        if shape.placeholder_format.idx == 0:
            continue
        if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            continue
            
        # Calculate area
        area = shape.width * shape.height
        content_placeholders.append((shape, area))
    
    if not content_placeholders:
        return None
    
    # Return the placeholder with the largest area
    return max(content_placeholders, key=lambda x: x[1])[0]

def get_placeholder_safely(slide, idx, placeholder_name=""):
    """
    Safely get a placeholder from a slide by index.
    Falls back to the largest non-title placeholder if index not found.
    
    Args:
        slide: The slide object
        idx: The placeholder index to retrieve
        placeholder_name: Optional name for better error messages
    
    Returns:
        The placeholder object if found, None otherwise
    """
    try:
        return slide.placeholders[idx]
    except KeyError:
        # Get all available placeholder indices
        available_indices = [ph.placeholder_format.idx for ph in slide.placeholders]
        
        logger.warning(
            f"Placeholder {idx} ({placeholder_name}) not found in slide. "
            f"Available indices: {available_indices}"
        )
        
        # Try to find the largest non-title placeholder
        largest_placeholder = get_largest_content_placeholder(slide)
        
        if largest_placeholder:
            logger.info(f"Using largest content placeholder at index {largest_placeholder.placeholder_format.idx}")
            return largest_placeholder
        
        # Last resort: try first available non-title placeholder
        non_title_indices = [idx for idx in available_indices if idx != 0]
        if non_title_indices:
            fallback_idx = non_title_indices[0]
            logger.info(f"Using fallback placeholder at index {fallback_idx}")
            return slide.placeholders[fallback_idx]
        
        logger.error("No placeholders found in this slide at all!")
        return None


def get_placeholder_by_type(slide, placeholder_type):
    """
    Get placeholder by type instead of index (more robust).
    
    Args:
        slide: The slide object
        placeholder_type: Type of placeholder (e.g., PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.BODY)
    
    Returns:
        The first placeholder of the specified type, or None
    """
    for shape in slide.placeholders:
        if shape.placeholder_format.type == placeholder_type:
            return shape
    return None

def remove_slide_number_from_heading(header: str) -> str:
    """
    Remove the slide number from a given slide header.

    Args:
        header: The header of a slide.

    Returns:
        str: The header without slide number.
    """
    if SLIDE_NUMBER_REGEX.match(header):
        idx = header.find(':')
        header = header[idx + 1:].strip()

    return header


def add_bulleted_items(text_frame: pptx.text.text.TextFrame, flat_items_list: list):
    """Add a list of texts as bullet points to a text frame and apply formatting.

    Args:
        text_frame (pptx.text.text.TextFrame): The text frame where text is to be
            displayed.
        flat_items_list (list): The list of items to be displayed.
    """

    for idx, an_item in enumerate(flat_items_list):
        if idx == 0:
            paragraph = text_frame.paragraphs[0]  # First paragraph for title text
        else:
            paragraph = text_frame.add_paragraph()
            paragraph.level = an_item[1]

        format_text(paragraph, an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER))


def format_text(frame_paragraph, text: str):
    """
    Apply bold and italic formatting while preserving the original word order without duplication.

    Args:
        frame_paragraph: The paragraph to format.
        text: The text to format with markdown-style formatting.
    """
    matches = list(BOLD_ITALICS_PATTERN.finditer(text))
    last_index = 0  # Track position in the text
    # Group 0: Full match (e.g., **bold** or *italic*)
    # Group 1: The outer parentheses (captures either bold or italic match, because of |)
    # Group 2: The bold text inside **bold**
    # Group 3: The italic text inside *italic*
    for match in matches:
        start, end = match.span()
        # Add unformatted text before the formatted section
        if start > last_index:
            run = frame_paragraph.add_run()
            run.text = text[last_index:start]

        # Extract formatted text
        if match.group(2):  # Bold
            run = frame_paragraph.add_run()
            run.text = match.group(2)
            run.font.bold = True
        elif match.group(3):  # Italics
            run = frame_paragraph.add_run()
            run.text = match.group(3)
            run.font.italic = True

        last_index = end  # Update position

    # Add any remaining unformatted text
    if last_index < len(text):
        run = frame_paragraph.add_run()
        run.text = text[last_index:]


def generate_powerpoint_presentation(
        parsed_data: dict,
        slides_template: str,
        output_file_path: pathlib.Path
) -> list:
    """
    Create and save a PowerPoint presentation from parsed JSON content.

    Args:
        parsed_data (dict): The presentation content as parsed JSON data.
        slides_template (str): The PPTX template key to use from GlobalConfig.
        output_file_path (pathlib.Path): Destination path for the generated PPTX file.

    Returns:
        A list containing the presentation title and slide headers.
    """

    presentation = pptx.Presentation(GlobalConfig.PPTX_TEMPLATE_FILES[slides_template]['file'])
    slide_id_list = [slide.slide_id for slide in presentation.slides]
    for slide_id in slide_id_list:
        rId = presentation.slides._sldIdLst[0].rId
        presentation.part.drop_rel(rId)
        del presentation.slides._sldIdLst[0]
    slide_width_inch, slide_height_inch = _get_slide_width_height_inches(presentation)

    # The title slide
    title_slide_layout = presentation.slide_layouts[0]
    
    slide = presentation.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    
    presentation_title = parsed_data.get('title', 'Untitled Presentation')
    
    try:
        if title is not None:
            title.text = presentation_title
        else:
            logger.warning("No title placeholder found in template slide")
            # Try to use subtitle as fallback
            if subtitle is not None:
                subtitle.text = presentation_title
            else:
                logger.error("No title or subtitle placeholder found - template incompatible")
    except Exception as e:
        logger.error(f"Error setting title: {e}")
        presentation_title = "Untitled Presentation"
    
    # Use presentation_title variable instead of title.text for logging
    logger.info(
        'PPT title: %s | #slides: %d | template: %s',
        presentation_title,  # Changed from title.text
        len(parsed_data['slides']),
        GlobalConfig.PPTX_TEMPLATE_FILES[slides_template]['file']
    )
    
    if subtitle is not None:
        subtitle.text = 'by RRD :)'
    
    all_headers = [presentation_title, ]  
    subtitle.text = 'by RRD :)'
    # all_headers = [title.text, ]

    # Add content in a loop
    for a_slide in parsed_data['slides']:
        try:
            is_processing_done = _handle_icons_ideas(
                presentation=presentation,
                slide_json=a_slide,
                slide_width_inch=slide_width_inch,
                slide_height_inch=slide_height_inch
            )

            if not is_processing_done:
                is_processing_done = _handle_table(
                    presentation=presentation,
                    slide_json=a_slide,
                    slide_width_inch=slide_width_inch,
                    slide_height_inch=slide_height_inch
                )

            if not is_processing_done:
                is_processing_done = _handle_double_col_layout(
                    presentation=presentation,
                    slide_json=a_slide,
                    slide_width_inch=slide_width_inch,
                    slide_height_inch=slide_height_inch
                )

            if not is_processing_done:
                is_processing_done = _handle_step_by_step_process(
                    presentation=presentation,
                    slide_json=a_slide,
                    slide_width_inch=slide_width_inch,
                    slide_height_inch=slide_height_inch
                )

            if not is_processing_done:
                _handle_default_display(
                    presentation=presentation,
                    slide_json=a_slide,
                    slide_width_inch=slide_width_inch,
                    slide_height_inch=slide_height_inch
            )

        except Exception:
            # In case of any unforeseen error, try to salvage what is available
            logger.error(
                'An error occurred while processing a slide...continuing with the next one',
                exc_info=True
            )
            continue

    # The thank-you slide
    last_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(last_slide_layout)
    title = slide.shapes.title
    if title is not None:
        title.text = 'Thank you!'
    else:
        logger.warning("No title placeholder for thank you slide")

    presentation.save(output_file_path)

    return all_headers


def get_flat_list_of_contents(items: list, level: int) -> list[tuple]:
    """
    Flatten a (hierarchical) list of bullet points to a single list containing each item and
     its level.

    Args:
        items: A bullet point (string or list).
        level: The current level of hierarchy.

    Returns:
        A list of (bullet item text, hierarchical level) tuples.
    """

    flat_list = []

    for item in items:
        if isinstance(item, str):
            flat_list.append((item, level))
        elif isinstance(item, list):
            flat_list = flat_list + get_flat_list_of_contents(item, level + 1)

    return flat_list


def get_slide_placeholders(
        slide: pptx.slide.Slide,
        layout_number: int,
        is_debug: bool = False
) -> list[tuple[int, str]]:
    """
    Return the index and name (lower case) of all placeholders present in a
    slide, except the title placeholder.

    A placeholder in a slide is a place to add content. Each placeholder has a
    name and an index. This index is not a list index; it is a key used to look up
    a dict and may be non-contiguous. The title placeholder always has index 0.
    User-added placeholders get indices starting from 10.

    With user-edited or added placeholders, indices may be difficult to track. This
    function returns the placeholders' names as well, which may help distinguish
    between placeholders.

    Args:
        slide: The slide.
        layout_number: The layout number used by the slide.
        is_debug: Whether to print debugging statements.

    Returns:
        list[tuple[int, str]]: A list of (index, name) tuples for placeholders
        present in the slide, excluding the title placeholder.
    """

    if is_debug:
        print(
            f'Slide layout #{layout_number}:'
            f' # of placeholders: {len(slide.shapes.placeholders)} (including the title)'
        )

    placeholders = [
        (shape.placeholder_format.idx, shape.name.lower()) for shape in slide.shapes.placeholders
    ]
    placeholders.pop(0)  # Remove the title placeholder

    if is_debug:
        print(placeholders)

    return placeholders


def _handle_default_display(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
):
    """
    Display a list of text in a slide.
    """
    status = False

    if 'img_keywords' in slide_json:
        if random.random() < IMAGE_DISPLAY_PROBABILITY:
            if random.random() < FOREGROUND_IMAGE_PROBABILITY:
                status = _handle_display_image__in_foreground(
                    presentation,
                    slide_json,
                    slide_width_inch,
                    slide_height_inch
                )
            else:
                status = _handle_display_image__in_background(
                    presentation,
                    slide_json,
                    slide_width_inch,
                    slide_height_inch
                )

    if status:
        return

    # Image display failed, so display only text
    bullet_slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(bullet_slide_layout)

    shapes = slide.shapes
    title_shape = shapes.title

    try:
        body_shape = shapes.placeholders[1]
    except KeyError:
        placeholders = get_slide_placeholders(slide, layout_number=1)
        body_shape = shapes.placeholders[placeholders[0][0]]

    if title_shape is not None:
        title_shape.text = remove_slide_number_from_heading(slide_json['heading'])
    else:
        logger.warning("No title placeholder in layout 1 (default display)")
        
    text_frame = body_shape.text_frame

    flat_items_list = get_flat_list_of_contents(slide_json['bullet_points'], level=0)
    add_bulleted_items(text_frame, flat_items_list)

    _handle_key_message(
        the_slide=slide,
        slide_json=slide_json,
        slide_height_inch=slide_height_inch,
        slide_width_inch=slide_width_inch
    )

    remove_unused_placeholders(slide)

def _handle_display_image__in_foreground(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """
    Create a slide with text and image using a picture placeholder layout. If not image keyword is
    available, it will add only text to the slide.

    Args:
        presentation: The presentation object.
        slide_json: The content of the slide as JSON data.
        slide_width_inch: The width of the slide in inches.
        slide_height_inch: The height of the slide in inches.

    Returns:
        bool: True if the side has been processed.
    """

    img_keywords = slide_json.get('img_keywords', '').strip()
    
    slide_layout = presentation.slide_layouts[8]  # Picture with Caption
    slide = presentation.slides.add_slide(slide_layout)
    placeholders = None

    title_placeholder = slide.shapes.title
    if title_placeholder is not None:
        title_placeholder.text = remove_slide_number_from_heading(slide_json['heading'])
    else:
        logger.warning("No title placeholder in layout 8 (Picture with Caption)")

    try:
        pic_col: PicturePlaceholder = slide.shapes.placeholders[1]
    except KeyError:
        placeholders = get_slide_placeholders(slide, layout_number=8)
        pic_col = None
        for idx, name in placeholders:
            if 'picture' in name:
                pic_col: PicturePlaceholder = slide.shapes.placeholders[idx]

    try:
        text_col: SlidePlaceholder = slide.shapes.placeholders[2]
    except KeyError:
        text_col = None
        if not placeholders:
            placeholders = get_slide_placeholders(slide, layout_number=8)

        for idx, name in placeholders:
            if 'content' in name:
                text_col: SlidePlaceholder = slide.shapes.placeholders[idx]

    if text_col and 'bullet_points' in slide_json:
        flat_items_list = get_flat_list_of_contents(slide_json['bullet_points'], level=0)
        add_bulleted_items(text_col.text_frame, flat_items_list)

    if img_keywords and pic_col:
        try:
            photo_url, page_url = ims.get_photo_url_from_api_response(
                ims.search_pexels(query=img_keywords, size='medium')
            )

            if photo_url:
                pic_col.insert_picture(
                    ims.get_image_from_url(photo_url)
                )

                _add_text_at_bottom(
                    slide=slide,
                    slide_width_inch=slide_width_inch,
                    slide_height_inch=slide_height_inch,
                    text='Photo provided by Pexels',
                    hyperlink=page_url
                )
        except Exception as ex:
            logger.error(
                '*** Error occurred while adding image to slide: %s',
                str(ex)
            )
    
    remove_unused_placeholders(slide)
    return True


def _handle_display_image__in_background(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """
    Add a slide with text and an image in the background. It works just like
    `_handle_default_display()` but with a background image added. If not image keyword is
    available, it will add only text to the slide.

    Args:
        presentation: The presentation object.
        slide_json: The content of the slide as JSON data.
        slide_width_inch: The width of the slide in inches.
        slide_height_inch: The height of the slide in inches.

    Returns:
        True if the slide has been processed.
    """

    img_keywords = slide_json.get('img_keywords', '').strip()

    # Add a photo in the background, text in the foreground
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title_shape = slide.shapes.title

    try:
        body_shape = slide.shapes.placeholders[1]
    except KeyError:
        placeholders = get_slide_placeholders(slide, layout_number=1)
        body_shape = slide.shapes.placeholders[placeholders[0][0]]

    if title_shape is not None:
        title_shape.text = remove_slide_number_from_heading(slide_json['heading'])
    else:
        logger.warning("No title placeholder in layout 1 (background image)")
    
    if 'bullet_points' in slide_json:
        flat_items_list = get_flat_list_of_contents(slide_json['bullet_points'], level=0)
        add_bulleted_items(body_shape.text_frame, flat_items_list)

    if img_keywords:
        try:
            photo_url, page_url = ims.get_photo_url_from_api_response(
                ims.search_pexels(query=img_keywords, size='large')
            )

            if photo_url:
                picture = slide.shapes.add_picture(
                    image_file=ims.get_image_from_url(photo_url),
                    left=0,
                    top=0,
                    width=pptx.util.Inches(slide_width_inch),
                )

                try:
                    blip_elements = picture._element.xpath('.//a:blip')
                    if not blip_elements:
                        logger.warning('No blip element found in the picture.')
                    else:
                        for blip in blip_elements:
                            alpha_mod = blip.makeelement(
                                '{http://schemas.openxmlformats.org/drawingml/2006/main}alphaModFix'
                            )
                            alpha_mod.set('amt', '50000')  # 50% opacity

                            existing_alpha_mod = blip.find(
                                '{http://schemas.openxmlformats.org/drawingml/2006/main}alphaModFix'
                            )
                            if existing_alpha_mod is not None:
                                blip.remove(existing_alpha_mod)

                            blip.append(alpha_mod)
                except Exception as ex:
                    logger.error('Failed to apply transparency: %s', str(ex))

                _add_text_at_bottom(
                    slide=slide,
                    slide_width_inch=slide_width_inch,
                    slide_height_inch=slide_height_inch,
                    text='Photo provided by Pexels',
                    hyperlink=page_url
                )

                # Move picture to background
                try:
                    slide.shapes._spTree.remove(picture._element)
                    slide.shapes._spTree.insert(2, picture._element)
                except Exception as ex:
                    logger.error('Failed to move image to background: %s', str(ex))

        except Exception as ex:
            logger.error('*** Error adding background image: %s', str(ex))

    remove_unused_placeholders(slide)
    return True


def _handle_icons_ideas(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
):
    """
    Add a slide with some icons and text.
    If no suitable icons are found, the step numbers are shown.

    Args:
        presentation: The presentation object.
        slide_json: The content of the slide as JSON data.
        slide_width_inch: The width of the slide in inches.
        slide_height_inch: The height of the slide in inches.

    Returns:
        True if the slide has been processed.
    """

    if 'bullet_points' in slide_json and slide_json['bullet_points']:
        items = slide_json['bullet_points']

        # Ensure that it is a single list of strings without any sub-list
        for step in items:
            if not isinstance(step, str) or not step.startswith(ICON_BEGINNING_MARKER):
                return False

        slide_layout = presentation.slide_layouts[5]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        if title is not None:
            title.text = remove_slide_number_from_heading(slide_json['heading'])
        else:
            logger.warning("No title placeholder in layout 5 (icons)")

        n_items = len(items)
        text_box_size = INCHES_2

        total_width = n_items * ICON_SIZE
        spacing = (pptx.util.Inches(slide_width_inch) - total_width) / (n_items + 1)
        top = INCHES_3

        icons_texts = [
            (match.group(1), match.group(2)) for match in [
                ICONS_REGEX.search(item) for item in items
            ]
        ]
        fallback_icon_files = ice.find_icons([item[0] for item in icons_texts])

        for idx, item in enumerate(icons_texts):
            icon, accompanying_text = item
            icon_path = f'{GlobalConfig.ICONS_DIR}/{icon}.png'

            if not os.path.exists(icon_path):
                logger.warning(
                    'Icon not found: %s...using fallback icon: %s',
                    icon, fallback_icon_files[idx]
                )
                icon_path = f'{GlobalConfig.ICONS_DIR}/{fallback_icon_files[idx]}.png'

            left = spacing + idx * (ICON_SIZE + spacing)
            center = left + ICON_SIZE / 2

            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                center - INCHES_0_5,
                top - (ICON_BG_SIZE - ICON_SIZE) / 2,
                INCHES_1, INCHES_1
            )
            shape.fill.solid()
            shape.shadow.inherit = False
            shape.fill.fore_color.rgb = shape.line.color.rgb = random.choice(ICON_COLORS)
            
            slide.shapes.add_picture(icon_path, left, top, height=ICON_SIZE)

            text_box = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                left=center - text_box_size / 2,
                top=top + ICON_SIZE + INCHES_0_2,
                width=text_box_size,
                height=text_box_size
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.CENTER
            format_text(text_frame.paragraphs[0], accompanying_text)

            text_frame.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
            text_box.fill.background()
            text_box.line.fill.background()
            text_box.shadow.inherit = False

            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.theme_color = pptx.enum.dml.MSO_THEME_COLOR.TEXT_2

        _add_text_at_bottom(
            slide=slide,
            slide_width_inch=slide_width_inch,
            slide_height_inch=slide_height_inch,
            text='More icons available in the SlideDeck AI repository',
            hyperlink='https://github.com/barun-saha/slide-deck-ai/tree/main/icons/png128'
        )

        remove_unused_placeholders(slide)
        return True
    
    return False


def _add_text_at_bottom(
        slide: pptx.slide.Slide,
        slide_width_inch: float,
        slide_height_inch: float,
        text: str,
        hyperlink: Optional[str] = None,
        target_height: Optional[float] = 0.5
):
    """
    Add arbitrary text to a textbox positioned near the lower-left side of a slide.

    Args:
        slide: The slide.
        slide_width_inch: The width of the slide in inches.
        slide_height_inch: The height of the slide in inches.
        text: The text to be added.
        hyperlink: Optional; the hyperlink to be added to the text.
        target_height: Optional[float]; the target height of the box in inches.
    """

    footer = slide.shapes.add_textbox(
        left=INCHES_1,
        top=pptx.util.Inches(slide_height_inch - target_height),
        width=pptx.util.Inches(slide_width_inch),
        height=pptx.util.Inches(target_height)
    )

    paragraph = footer.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = pptx.util.Pt(10)
    run.font.underline = False

    if hyperlink:
        run.hyperlink.address = hyperlink

def _handle_double_col_layout(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """
    Add a slide with a double column layout for comparison.

    Args:
        presentation (pptx.Presentation): The presentation object.
        slide_json (dict): The content of the slide as JSON data.
        slide_width_inch (float): The width of the slide in inches.
        slide_height_inch (float): The height of the slide in inches.

    Returns:
        bool: True if double col layout has been added; False otherwise.
    """

    if 'bullet_points' in slide_json and slide_json['bullet_points']:
        double_col_content = slide_json['bullet_points']

        if double_col_content and (
                len(double_col_content) == 2
        ) and isinstance(double_col_content[0], dict) and isinstance(double_col_content[1], dict):
            slide_layout = presentation.slide_layouts[4]
            slide = presentation.slides.add_slide(slide_layout)
            placeholders = None

            shapes = slide.shapes
            title_placeholder = shapes.title
            if title_placeholder is not None:
                title_placeholder.text = remove_slide_number_from_heading(slide_json['heading'])
            else:
                logger.warning("No title placeholder in layout 4 (double column)")

            try:
                left_heading, right_heading = shapes.placeholders[1], shapes.placeholders[3]
            except KeyError:
                left_heading, right_heading = None, None
                placeholders = get_slide_placeholders(slide, layout_number=4)

                for idx, name in placeholders:
                    if 'text placeholder' in name:
                        if not left_heading:
                            left_heading = shapes.placeholders[idx]
                        elif not right_heading:
                            right_heading = shapes.placeholders[idx]

            try:
                left_col, right_col = shapes.placeholders[2], shapes.placeholders[4]
            except KeyError:
                left_col, right_col = None, None
                if not placeholders:
                    placeholders = get_slide_placeholders(slide, layout_number=4)

                for idx, name in placeholders:
                    if 'content placeholder' in name:
                        if not left_col:
                            left_col = shapes.placeholders[idx]
                        elif not right_col:
                            right_col = shapes.placeholders[idx]

            if left_col and right_col:
                left_col_frame, right_col_frame = left_col.text_frame, right_col.text_frame

                if 'heading' in double_col_content[0] and left_heading:
                    left_heading.text = double_col_content[0]['heading']
                if 'bullet_points' in double_col_content[0]:
                    flat_items_list = get_flat_list_of_contents(
                        double_col_content[0]['bullet_points'], level=0
                    )

                    if not left_heading:
                        left_col_frame.text = double_col_content[0]['heading']

                    add_bulleted_items(left_col_frame, flat_items_list)

                if 'heading' in double_col_content[1] and right_heading:
                    right_heading.text = double_col_content[1]['heading']
                if 'bullet_points' in double_col_content[1]:
                    flat_items_list = get_flat_list_of_contents(
                        double_col_content[1]['bullet_points'], level=0
                    )

                    if not right_heading:
                        right_col_frame.text = double_col_content[1]['heading']

                    add_bulleted_items(right_col_frame, flat_items_list)

                _handle_key_message(
                    the_slide=slide,
                    slide_json=slide_json,
                    slide_height_inch=slide_height_inch,
                    slide_width_inch=slide_width_inch
                )

                remove_unused_placeholders(slide)
                return True
    
    return False


def _handle_step_by_step_process(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """Add shapes to display a step-by-step process in the slide, if available.

    Args:
        presentation (pptx.Presentation): The presentation object.
        slide_json (dict): The content of the slide as JSON data.
        slide_width_inch (float): The width of the slide in inches.
        slide_height_inch (float): The height of the slide in inches.

    Returns:
        bool: True if this slide has a step-by-step process depiction added; False otherwise.
    """

    if 'bullet_points' in slide_json and slide_json['bullet_points']:
        steps = slide_json['bullet_points']

        no_marker_count = 0.0
        n_steps = len(steps)

        for step in steps:
            if not isinstance(step, str):
                return False

            if not step.startswith(STEP_BY_STEP_PROCESS_MARKER):
                no_marker_count += 1

        slide_header = slide_json['heading'].lower()
        if (no_marker_count / n_steps > 0.25) and not (
                ('step-by-step' in slide_header) or ('step by step' in slide_header)
        ):
            return False

        if n_steps < 3 or n_steps > 6:
            return False

        bullet_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title = shapes.title
        if title is not None:
            title.text = remove_slide_number_from_heading(slide_json['heading'])
        else:
            logger.warning("No title placeholder in step-by-step layout")

        if 3 <= n_steps <= 4:
            # Horizontal display
            height = INCHES_1_5
            width = pptx.util.Inches(slide_width_inch / n_steps - 0.01)
            top = pptx.util.Inches(slide_height_inch / 2)
            left = pptx.util.Inches((slide_width_inch - width.inches * n_steps) / 2 + 0.05)

            for step in steps:
                shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
                text_frame = shape.text_frame
                text_frame.clear()
                paragraph = text_frame.paragraphs[0]
                paragraph.alignment = pptx.enum.text.PP_ALIGN.LEFT
                format_text(paragraph, step.removeprefix(STEP_BY_STEP_PROCESS_MARKER))
                left += width - INCHES_0_4
        elif 4 < n_steps <= 6:
            # Vertical display
            height = pptx.util.Inches(0.65)
            top = pptx.util.Inches(slide_height_inch / 4)
            left = INCHES_1

            width = pptx.util.Inches(slide_width_inch * 2 / 3)
            lengths = [len(step) for step in steps]
            font_size_20pt = pptx.util.Pt(20)
            widths = sorted(
                [
                    min(
                        pptx.util.Inches(font_size_20pt.inches * a_len),
                        width
                    ) for a_len in lengths
                ]
            )
            width = widths[len(widths) // 2]

            for step in steps:
                shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON, left, top, width, height)
                text_frame = shape.text_frame
                text_frame.clear()
                paragraph = text_frame.paragraphs[0]
                paragraph.alignment = pptx.enum.text.PP_ALIGN.LEFT
                format_text(paragraph, step.removeprefix(STEP_BY_STEP_PROCESS_MARKER))
                top += height + INCHES_0_3
                left += INCHES_0_5

        remove_unused_placeholders(slide)
        return True

    return False


def _handle_table(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """
    Add a table to a slide, if available.
    """
    if 'table' not in slide_json or not slide_json['table']:
        return False
    
    headers = slide_json['table'].get('headers', [])
    rows = slide_json['table'].get('rows', [])
    
    if not headers or not rows:
        return False
    
    # Create the slide FIRST
    bullet_slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    # Set the title
    title = shapes.title
    if title is not None:
        title.text = remove_slide_number_from_heading(slide_json['heading'])
    else:
        logger.warning("No title placeholder in table layout")
    
    # Try to get the largest content placeholder
    placeholder = get_largest_content_placeholder(slide)
    
    # Determine table position and size
    if placeholder is None:
        logger.warning("No placeholder found for table - using default positioning")
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(slide_width_inch - 1)
        height = Inches(slide_height_inch - 3)
    else:
        left = placeholder.left
        top = placeholder.top
        width = placeholder.width
        height = placeholder.height
    
    # Create the table
    table = slide.shapes.add_table(
        len(rows) + 1,  # rows + 1 for header
        len(headers),   # columns
        left, 
        top, 
        width, 
        height
    ).table
    
    # Set headers
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header_text
        cell.text_frame.paragraphs[0].font.bold = True
    
    # Fill in rows
    for row_idx, row_data in enumerate(rows, start=1):
        for col_idx, cell_text in enumerate(row_data):
            table.cell(row_idx, col_idx).text = str(cell_text)
    
    remove_unused_placeholders(slide)
    return True

def _handle_key_message(
        the_slide: pptx.slide.Slide,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
):
    """
        Add a shape to display the key message in the slide, if available.

        Args:
            the_slide (pptx.slide.Slide): The slide to be processed.
            slide_json (dict): The content of the slide as JSON data.
            slide_width_inch (float): The width of the slide in inches.
            slide_height_inch (float): The height of the slide in inches.

        Returns:
            None
        """

    if 'key_message' in slide_json and slide_json['key_message']:
        height = pptx.util.Inches(1.6)
        width = pptx.util.Inches(slide_width_inch / 2.3)
        top = pptx.util.Inches(slide_height_inch - height.inches - 0.1)
        left = pptx.util.Inches((slide_width_inch - width.inches) / 2)
        shape = the_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left=left,
            top=top,
            width=width,
            height=height
        )
        format_text(shape.text_frame.paragraphs[0], slide_json['key_message'])


def _get_slide_width_height_inches(presentation: pptx.Presentation) -> tuple[float, float]:
    """
    Get the dimensions of a slide in inches.

    Args:
        presentation: The presentation object.

    Returns:
        The width and the height.
    """

    slide_width_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_width
    slide_height_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_height

    return slide_width_inch, slide_height_inch


if __name__ == '__main__':
    _JSON_DATA = '''
    {
  "title": "AI Applications: Transforming Industries",
  "slides": [
    {
      "heading": "Introduction to AI Applications",
      "bullet_points": [
        "Artificial Intelligence (AI) is *transforming* various industries",
        "AI applications range from simple decision-making tools to complex systems",
        "AI can be categorized into types: Rule-based, Instance-based, and Model-based"
      ],
      "key_message": "AI is a broad field with diverse applications and categories",
      "img_keywords": "AI, transformation, industries, decision-making, categories"
    },
    {
      "heading": "AI in Everyday Life",
      "bullet_points": [
        "**Virtual assistants** like Siri, Alexa, and Google Assistant",
        "**Recommender systems** in Netflix, Amazon, and Spotify",
        "**Fraud detection** in banking and *credit card* transactions"
      ],
      "key_message": "AI is integrated into our daily lives through various services",
      "img_keywords": "virtual assistants, recommender systems, fraud detection"
    },
    {
      "heading": "AI in Healthcare",
      "bullet_points": [
        "Disease diagnosis and prediction using machine learning algorithms",
        "Personalized medicine and drug discovery",
        "AI-powered robotic surgeries and remote patient monitoring"
      ],
      "key_message": "AI is revolutionizing healthcare with improved diagnostics and patient care",
      "img_keywords": "healthcare, disease diagnosis, personalized medicine, robotic surgeries"
    },
    {
      "heading": "AI in Key Industries",
      "bullet_points": [
        {
          "heading": "Retail",
          "bullet_points": [
            "Inventory management and demand forecasting",
            "Customer segmentation and targeted marketing",
            "AI-driven chatbots for customer service"
          ]
        },
        {
          "heading": "Finance",
          "bullet_points": [
            "Credit scoring and risk assessment",
            "Algorithmic trading and portfolio management",
            "AI for detecting money laundering and cyber fraud"
          ]
        }
      ],
      "key_message": "AI is transforming retail and finance with improved operations and decision-making",
      "img_keywords": "retail, finance, inventory management, credit scoring, algorithmic trading"
    },
    {
      "heading": "AI in Education",
      "bullet_points": [
        "Personalized learning paths and adaptive testing",
        "Intelligent tutoring systems for skill development",
        "AI for predicting student performance and dropout rates"
      ],
      "key_message": "AI is personalizing education and improving student outcomes",
    },
    {
      "heading": "Step-by-Step: AI Development Process",
      "bullet_points": [
        ">> **Step 1:** Define the problem and objectives",
        ">> **Step 2:** Collect and preprocess data",
        ">> **Step 3:** Select and train the AI model",
        ">> **Step 4:** Evaluate and optimize the model",
        ">> **Step 5:** Deploy and monitor the AI system"
      ],
      "key_message": "Developing AI involves a structured process from problem definition to deployment",
      "img_keywords": ""
    },
    {
      "heading": "AI Icons: Key Aspects",
      "bullet_points": [
        "[[brain]] Human-like *intelligence* and decision-making",
        "[[robot]] Automation and physical *tasks*",
        "[[]] Data processing and cloud computing",
        "[[lightbulb]] Insights and *predictions*",
        "[[globe2]] Global connectivity and *impact*"
      ],
      "key_message": "AI encompasses various aspects, from human-like intelligence to global impact",
      "img_keywords": "AI aspects, intelligence, automation, data processing, global impact"
    },
    {
        "heading": "AI vs. ML vs. DL: A Tabular Comparison",
        "table": {
            "headers": ["Feature", "AI", "ML", "DL"],
            "rows": [
                ["Definition", "Creating intelligent agents", "Learning from data", "Deep neural networks"],
                ["Approach", "Rule-based, expert systems", "Algorithms, statistical models", "Deep neural networks"],
                ["Data Requirements", "Varies", "Large datasets", "Massive datasets"],
                ["Complexity", "Varies", "Moderate", "High"],
                ["Computational Cost", "Low to Moderate", "Moderate", "High"],
                ["Examples", "Chess, recommendation systems", "Spam filters, image recognition", "Image recognition, NLP"]
            ]
        },
        "key_message": "This table provides a concise comparison of the key features of AI, ML, and DL.",
        "img_keywords": "AI, ML, DL, comparison, table, features"
    },
    {
      "heading": "Conclusion: Embracing AI's Potential",
      "bullet_points": [
        "AI is transforming industries and improving lives",
        "Ethical considerations are crucial for responsible AI development",
        "Invest in AI education and workforce development",
        "Call to action: Explore AI applications and contribute to shaping its future"
      ],
      "key_message": "AI offers *immense potential*, and we must embrace it responsibly",
      "img_keywords": "AI transformation, ethical considerations, AI education, future of AI"
    }
  ]
}'''

    temp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    path = pathlib.Path(temp.name)

    generate_powerpoint_presentation(
        json5.loads(_JSON_DATA),
        output_file_path=path,
        slides_template='Basic'
    )
    print(f'File path: {path}')

    temp.close()

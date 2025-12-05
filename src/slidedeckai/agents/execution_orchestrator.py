# slidedeckai/agents/execution_orchestrator.py - FULLY FIXED
"""
Orchestrates the execution of a slide deck generation plan.

This module coordinates the content generation, search execution, and PowerPoint
slide creation. It extracts template properties, manages slide layouts,
fills placeholders with appropriate content (text, charts, tables, icons),
and handles the final presentation assembly.
"""
import logging
import pathlib
import json
from typing import Dict, List, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from concurrent.futures import ThreadPoolExecutor, as_completed

from .search_executor import WebSearchExecutor
from .content_generator import ContentGenerator
from slidedeckai.layout_analyzer import TemplateAnalyzer
from slidedeckai.content_matcher import ContentLayoutMatcher
from slidedeckai.helpers.icon_selector import IconSelector
from openai import OpenAI

logger = logging.getLogger(__name__)


class ExecutionOrchestrator:
    """
    Executes a research plan to create a PowerPoint presentation.

    This class handles the end-to-end process of generating slides, including:
    - Parallel web searches.
    - Content generation using LLMs.
    - Intelligent placeholder filling (text, charts, tables).
    - Template property extraction and application (fonts, colors).
    """
    
    def __init__(self, api_key: str, template_path: pathlib.Path, use_llm_role_validation: bool = False):
        """
        Initialize the orchestrator.

        Args:
            api_key (str): OpenAI API key.
            template_path (pathlib.Path): Path to the PowerPoint template file.
            use_llm_role_validation (bool): Whether to use LLM for advanced role validation of placeholders.
        """
        self.api_key = api_key
        self.template_path = template_path
        self.search_executor = WebSearchExecutor(api_key)
        self.content_generator = ContentGenerator(api_key)
        self.icon_selector = IconSelector()
        self.openai_client = OpenAI(api_key=api_key)
        # Optional: use the LLM to validate/override inferred placeholder roles
        self.use_llm_role_validation = use_llm_role_validation
        
        # Load template and extract properties
        self.presentation = Presentation(template_path)
        self.template_properties = self._extract_template_properties()
        # Analyzer & matcher for intelligent layout/content mapping
        try:
            self.analyzer = TemplateAnalyzer(self.presentation)
            self.matcher = ContentLayoutMatcher(self.analyzer)
        except Exception as e:
            logger.debug(f"Could not initialize TemplateAnalyzer/Matcher: {e}")
            self.analyzer = None
            self.matcher = None
        
    def _extract_template_properties(self) -> Dict:
        """
        Extract template properties (colors, fonts, spacing) dynamically.

        Returns:
            Dict: A dictionary containing extracted properties.
        """
        properties = {
            'slide_width': self.presentation.slide_width,
            'slide_height': self.presentation.slide_height,
            'theme_colors': {},
            'default_fonts': {'name': 'Calibri', 'size': Pt(18)},  # Default fallback
            'spacing': {'margin_left': 0.5, 'margin_top': 1.0, 'line_spacing': 1.5},
        }
        
        # Extract theme colors from first slide master
        try:
            if self.presentation.slide_master and hasattr(self.presentation.slide_master, 'theme'):
                theme = self.presentation.slide_master.theme
                if hasattr(theme, 'theme_colors') and len(theme.theme_colors) >= 6:
                    properties['theme_colors'] = {
                        'accent1': theme.theme_colors[4],  # Accent 1
                        'accent2': theme.theme_colors[5],  # Accent 2
                        'text1': theme.theme_colors[0] if len(theme.theme_colors) > 0 else RGBColor(0, 0, 0),
                        'text2': theme.theme_colors[1] if len(theme.theme_colors) > 1 else RGBColor(68, 68, 68),
                        'background1': theme.theme_colors[2] if len(theme.theme_colors) > 2 else RGBColor(255, 255, 255),
                    }
        except Exception as e:
            logger.debug(f"Could not extract theme colors: {e}")
        
        # Set default theme colors if extraction failed
        if not properties['theme_colors']:
            properties['theme_colors'] = {
                'accent1': RGBColor(68, 114, 196),
                'accent2': RGBColor(112, 173, 71),
                'text1': RGBColor(0, 0, 0),
                'text2': RGBColor(68, 68, 68),
                'background1': RGBColor(255, 255, 255),
            }
        
        # Extract font defaults from first layout with text
        try:
            for layout in self.presentation.slide_layouts:
                for shape in layout.placeholders:
                    try:
                        if shape.has_text_frame and shape.text_frame.paragraphs:
                            para = shape.text_frame.paragraphs[0]
                            if para.runs:
                                run = para.runs[0]
                                font_name = run.font.name or 'Calibri'
                                font_size = run.font.size or Pt(18)
                                properties['default_fonts'] = {'name': font_name, 'size': font_size}
                                break
                    except Exception:
                        continue
                if properties['default_fonts'].get('name') != 'Calibri' or properties['default_fonts'].get('size') != Pt(18):
                    break
        except Exception as e:
            logger.debug(f"Could not extract fonts: {e}")
        
        # Extract spacing from largest placeholder
        try:
            max_area = 0
            for layout in self.presentation.slide_layouts:
                for shape in layout.placeholders:
                    area = (shape.width / 914400.0) * (shape.height / 914400.0)
                    if area > max_area:
                        max_area = area
                        properties['spacing'] = {
                            'margin_left': shape.left / 914400.0,
                            'margin_top': shape.top / 914400.0,
                            'line_spacing': 1.5,  # Default
                        }
        except Exception as e:
            logger.warning(f"Could not extract spacing: {e}")
            properties['spacing'] = {
                'margin_left': 0.5,
                'margin_top': 1.0,
                'line_spacing': 1.5,
            }
        
        logger.info(f"✅ Extracted template properties: {len(properties['theme_colors'])} colors")
        return properties
    
    def execute_plan(self, plan: Any, output_path: pathlib.Path, chart_data: Optional[Dict] = None, extracted_content: Optional[str] = None) -> pathlib.Path:
        """
        Execute the research plan and generate the presentation.

        Args:
            plan (ResearchPlan): The plan to execute.
            output_path (pathlib.Path): The path to save the generated PPTX file.
            chart_data (Optional[Dict]): Optional chart data to override generation.
            extracted_content (Optional[str]): Optional content extracted from uploaded files.

        Returns:
            pathlib.Path: The path to the generated presentation.
        """
        # DEMO MODE SHORTCUT
        if plan.search_mode == "demo":
            logger.info("🤖 DEMO MODE: Generating mock presentation without LLM/Search")
            return self._execute_mock_plan(plan, output_path)

        logger.info("🚀 Executing FULLY FIXED plan...")
        logger.info(f"  Slides: {len(plan.sections)}")
        
        # STEP 1: Execute searches IN PARALLEL
        all_queries = []
        # If expected_source_type is 'extracted_content', we skip web search
        search_queries = []

        for section in plan.sections:
            for spec in section.placeholder_specs:
                for q in spec.search_queries:
                    if getattr(q, 'expected_source_type', '') != 'extracted_content':
                         search_queries.append(q.query)
        
        logger.info(f"  Queries: {len(search_queries)}")
        
        if search_queries:
            logger.info("🔍 Executing searches IN PARALLEL...")
            search_results = self._execute_searches_parallel(search_queries)
            logger.info(f"✅ {len(search_results)} searches complete")
        else:
            search_results = {}

        # If we have extracted content, make it available for content generation
        # by treating it as a "fact" for queries tagged with 'extracted_content'
        if extracted_content:
             # Iterate again to populate search_results with extracted_content
             for section in plan.sections:
                for spec in section.placeholder_specs:
                    for q in spec.search_queries:
                        if getattr(q, 'expected_source_type', '') == 'extracted_content':
                             # Use the extracted content as the result
                             # We truncate it slightly if it's too huge, but ideally we should search IN it.
                             # For now, we pass it all as one "fact"
                             search_results[q.query] = [extracted_content]
        
        # STEP 2: Clear existing slides (keep only master)
        slide_ids = [slide.slide_id for slide in self.presentation.slides]
        for slide_id in slide_ids:
            rId = self.presentation.slides._sldIdLst[0].rId
            self.presentation.part.drop_rel(rId)
            del self.presentation.slides._sldIdLst[0]
        
        # STEP 3: ADD TITLE SLIDE (FIX #2)
        logger.info("📄 Adding title slide...")
        self._add_title_slide(plan.query)
        
        # STEP 4: Generate content slides
        execution_log = []
        
        for idx, section in enumerate(plan.sections, 1):
            try:
                slide_log = self._generate_slide_smart(
                    section, 
                    search_results,
                    idx,
                    len(plan.sections),
                    chart_data=chart_data
                )
                execution_log.append(slide_log)
                
            except Exception as e:
                logger.error(f"❌ Slide {idx} failed: {e}", exc_info=True)
                execution_log.append({
                    'slide': idx,
                    'title': section.section_title,
                    'status': 'failed',
                    'error': str(e)
                })
        
        # STEP 5: ADD THANK YOU SLIDE (FIX #2)
        logger.info("📄 Adding thank you slide...")
        self._add_thank_you_slide()
        
        # STEP 6: Save
        self.presentation.save(output_path)
        logger.info(f"✅ Saved: {output_path}")
        
        # Save execution log
        log_path = str(output_path).replace('.pptx', '.execution.json')
        with open(log_path, 'w') as f:
            json.dump(execution_log, f, indent=2)
        logger.info(f"📋 Execution log saved: {log_path}")
        
        return output_path
    
    def _execute_searches_parallel(self, queries: List[str]) -> Dict[str, List[str]]:
        """
        Execute web searches in parallel using ThreadPoolExecutor.

        Args:
            queries (List[str]): List of search queries.

        Returns:
            Dict[str, List[str]]: Dictionary mapping queries to list of facts.
        """
        results = {}
        
        with ThreadPoolExecutor(max_workers=5) as executor:
            # Submit all searches
            future_to_query = {
                executor.submit(self.search_executor._search_with_gpt, query): query
                for query in queries
            }
            
            # Collect results as they complete
            for future in as_completed(future_to_query):
                query = future_to_query[future]
                try:
                    facts = future.result()
                    results[query] = facts
                    logger.info(f"  ✓ {query}: {len(facts)} facts")
                except Exception as e:
                    logger.error(f"  ✗ {query} failed: {e}")
                    results[query] = [f"Data for {query}: See latest reports"]
        
        return results

    def _prepare_section_content(self, section: Any, placeholder_map: Dict, search_results: Dict) -> Dict:
        """
        Generate content for all placeholders in a section in parallel.

        Args:
            section (Any): The section plan object.
            placeholder_map (Dict): Map of placeholder IDs to their info.
            search_results (Dict): Search results.

        Returns:
            Dict: Mapping of placeholder IDs to generated content.
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed
        results = {}

        def _gen_for_ph(ph_id, ph_info):
            role = ph_info.get('role')
            # Gather relevant facts
            relevant_facts = []
            for spec in getattr(section, 'placeholder_specs', []):
                try:
                    spec_idx = getattr(spec, 'placeholder_idx', None)
                    if spec_idx is None:
                        continue
                    try:
                        if int(spec_idx) != int(ph_id):
                            continue
                    except (ValueError, TypeError):
                        if spec_idx != ph_id:
                            continue
                        for q in getattr(spec, 'search_queries', []):
                            qq = getattr(q, 'query', None)
                            if qq in search_results:
                                relevant_facts.extend(search_results[qq])
                except Exception:
                    continue

            try:
                if role == 'subtitle':
                    text = self.content_generator.generate_subtitle(
                        section.section_title,
                        section.section_purpose,
                        relevant_facts[:3]
                    )
                    return (ph_id, {'type': 'subtitle', 'text': text})
                elif role == 'chart':
                    chart_data = self.content_generator.generate_chart(
                        section.section_title,
                        section.section_purpose,
                        relevant_facts,
                        chart_type='column'
                    )
                    return (ph_id, {'type': 'chart', 'chart_data': chart_data})
                elif role == 'table':
                    table_data = self.content_generator.generate_table(
                        section.section_title,
                        section.section_purpose,
                        relevant_facts
                    )
                    return (ph_id, {'type': 'table', 'table_data': table_data})
                elif role == 'kpi':
                    kpi = self.content_generator.generate_kpi(
                        section.section_title,
                        relevant_facts[0] if relevant_facts else f"KPI for {section.section_title}"
                    )
                    return (ph_id, {'type': 'kpi', 'kpi_data': kpi})
                else:
                    bullets = self.content_generator.generate_bullets(
                        section.section_title,
                        section.section_purpose,
                        relevant_facts,
                        max_bullets=self._calculate_max_bullets(ph_info.get('area', 5))
                    )
                    return (ph_id, {'type': 'bullets', 'bullets': bullets})
            except Exception as e:
                logger.error(f"Content generation failed for placeholder {ph_id}: {e}")
                return (ph_id, {'status': 'failed', 'error': str(e)})

        with ThreadPoolExecutor(max_workers=4) as executor:
            futures = {executor.submit(_gen_for_ph, ph_id, ph_info): ph_id for ph_id, ph_info in placeholder_map.items()}
            for future in as_completed(futures):
                ph_id = futures[future]
                try:
                    pid, data = future.result()
                    results[pid] = data
                except Exception as e:
                    logger.error(f"Placeholder {ph_id} generation exception: {e}")
                    results[ph_id] = {'status': 'failed', 'error': str(e)}

        return results
    
    def _add_title_slide(self, title: str):
        """
        Add a title slide to the presentation.

        Args:
            title (str): The title text.
        """
        title_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(title_layout)
        
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
        
        # Add subtitle if exists
        for shape in slide.placeholders:
            if shape.placeholder_format.idx != 0:
                try:
                    shape.text = 'Generated by SlideDeck AI'
                except:
                    pass
                break
        
        logger.info(f"  ✓ Title: {title}")
    
    def _add_thank_you_slide(self):
        """
        Add a 'Thank You' slide at the end of the presentation.
        """
        title_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(title_layout)
        
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = 'Thank You!'
        
        logger.info(f"  ✓ Thank you slide added")
    
    def _generate_slide_smart(self, section: Any, search_results: Dict,
                              slide_num: int, total: int, chart_data: Optional[Dict] = None) -> Dict:
        """
        Generate a single slide using smart layout analysis and content generation.

        Args:
            section (Any): The section plan.
            search_results (Dict): The search results.
            slide_num (int): The slide number.
            total (int): Total number of slides.
            chart_data (Optional[Dict]): Optional chart data.

        Returns:
            Dict: A log of the slide generation process.
        """
        
        layout_idx = section.layout_idx
        
        if not isinstance(layout_idx, int):
            layout_idx = int(layout_idx)
        
        logger.info(f"📄 Slide {slide_num}: {section.section_title} ({section.enforced_content_type})")
        logger.info(f"  Layout {layout_idx}: {section.layout_type}")
        
        # Get layout
        layout = self.presentation.slide_layouts[layout_idx]
        slide = self.presentation.slides.add_slide(layout)
        
        # Analyze placeholders
        placeholder_map = self._analyze_layout_placeholders(slide, layout_idx)

        # Integrate ContentLayoutMatcher suggestions (Gap 1 fix)
        try:
            if self.matcher and self.analyzer:
                layout_capability = None
                try:
                    layout_capability = self.analyzer.layouts.get(int(layout_idx))
                except Exception:
                    layout_capability = None

                # Build a minimal slide_json for matcher
                slide_json = {
                    'heading': getattr(section, 'section_title', ''),
                    'section_purpose': getattr(section, 'section_purpose', ''),
                    'bullet_points': []
                }
                # Populate bullets from placeholder_specs descriptions when available
                for spec in getattr(section, 'placeholder_specs', []) or []:
                    try:
                        desc = getattr(spec, 'content_description', None) or getattr(spec, 'content_type', None)
                        if desc:
                            slide_json['bullet_points'].append(str(desc))
                    except Exception:
                        continue

                if layout_capability:
                    try:
                        content_map = self.matcher.map_content_to_placeholders(slide_json, layout_capability)
                        # Merge suggestions into placeholder_map
                        for pid, spec in content_map.items():
                            try:
                                pid_key = int(pid) if isinstance(pid, (str, float)) else pid
                            except Exception:
                                pid_key = pid
                            if pid_key in placeholder_map and isinstance(spec, dict):
                                suggested_type = spec.get('type') or spec.get('role')
                                if suggested_type:
                                    placeholder_map[pid_key]['role'] = suggested_type
                                placeholder_map[pid_key]['suggested_content'] = spec
                    except Exception as e:
                        logger.debug(f"ContentLayoutMatcher mapping failed: {e}")
        except Exception:
            pass

        # PREPARE content for placeholders in parallel (only text/chart/table data generation)
        # If chart_data is provided globally, we inject it into prepared_content for chart placeholders
        prepared_content = self._prepare_section_content(section, placeholder_map, search_results)
        
        if chart_data:
             for ph_id, ph_info in placeholder_map.items():
                if ph_info['role'] == 'chart':
                    # Override/Inject chart data
                    prepared_content[ph_id] = {'type': 'chart', 'chart_data': chart_data}
                    logger.info(f"    ↳ Injected uploaded chart data for PH {ph_id}")

        logger.info(f"  📋 Layout has {len(placeholder_map)} placeholders:")
        for ph_id, ph_info in placeholder_map.items():
            logger.info(f"    [{ph_id}] {ph_info['type']} - {ph_info['area']:.1f} sq in - {ph_info['role']}")

        # Optional LLM-assisted role validation/override (batched)
        if getattr(self, 'use_llm_role_validation', False):
            logger.info("  🤖 Validating placeholder roles with LLM (batched)...")
            try:
                overrides = self._batch_validate_placeholder_roles(section, placeholder_map)
                for pid, new_role in overrides.items():
                    # Handle both int and str keys in placeholder_map
                    try:
                        pid_key = int(pid) if isinstance(pid, str) else pid
                    except (ValueError, TypeError):
                        pid_key = pid
                    
                    if pid_key in placeholder_map and new_role and new_role != placeholder_map[pid_key].get('role'):
                        old_role = placeholder_map[pid_key].get('role', 'unknown')
                        logger.info(f"    → Role override for ph {pid_key}: {old_role} -> {new_role}")
                        placeholder_map[pid_key]['role'] = new_role
            except Exception as e:
                logger.debug(f"Batched LLM role validation failed: {e}")
        
        # Set title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = section.section_title
            logger.info(f"    ✓ Title set")
        
        # Generate content for EACH placeholder
        slide_log = {
            'slide': slide_num,
            'title': section.section_title,
            'layout_idx': layout_idx,
            'layout_type': section.layout_type,
            'placeholders_found': len(placeholder_map),
            'placeholders': []
        }
        
        for ph_id, ph_info in placeholder_map.items():
            try:
                ph_log = self._fill_placeholder_smart(
                    slide,
                    ph_id,
                    ph_info,
                    section,
                    search_results,
                    prepared_content=prepared_content
                )
                slide_log['placeholders'].append(ph_log)
                
            except Exception as e:
                logger.error(f"    ❌ Placeholder {ph_id} failed: {e}")
                slide_log['placeholders'].append({
                    'id': ph_id,
                    'status': 'failed',
                    'error': str(e)
                })
        
        logger.info(f"  ✅ Complete")
        return slide_log
    
    def _analyze_layout_placeholders(self, slide: Any, layout_idx: int) -> Dict:
        """
        Analyze placeholders on a slide layout.

        Args:
            slide (Any): The slide object.
            layout_idx (int): Layout index.

        Returns:
            Dict: Analysis of placeholders.
        """
        
        placeholder_map = {}
        
        for shape in slide.placeholders:
            ph_idx = shape.placeholder_format.idx
            
            if ph_idx == 0:
                continue
            
            ph_type_id = shape.placeholder_format.type
            ph_type_name = self._get_placeholder_type_name(ph_type_id)
            
            try:
                left = shape.left / 914400.0 if shape.left else 0.0
                top = shape.top / 914400.0 if shape.top else 0.0
                width = shape.width / 914400.0 if shape.width else 1.0
                height = shape.height / 914400.0 if shape.height else 1.0
            except (AttributeError, TypeError, ZeroDivisionError):
                left, top, width, height = 0.0, 0.0, 1.0, 1.0
            area = width * height
            
            role = self._determine_placeholder_role(
                ph_type_id, ph_type_name, width, height, area
            )
            
            placeholder_map[ph_idx] = {
                'type': ph_type_name,
                'type_id': ph_type_id,
                'role': role,
                'left': left,
                'top': top,
                'width': width,
                'height': height,
                'area': area,
                'bbox': (left, top, width, height)
            }
        
        return placeholder_map
    
    def _determine_placeholder_role(self, type_id: int, type_name: str,
                                     width: float, height: float, area: float) -> str:
        """
        Determine the role of a placeholder based on its type and dimensions.

        Args:
            type_id (int): Placeholder type ID.
            type_name (str): Placeholder type name.
            width (float): Width in inches.
            height (float): Height in inches.
            area (float): Area in square inches.

        Returns:
            str: The role (e.g., 'subtitle', 'chart', 'content').
        """
        
        if type_id in [1, 4]:
            return 'subtitle'
        
        if type_id == 10:
            return 'chart'
        if type_id == 11:
            return 'table'
        if type_id == 15:
            return 'image'
        
        if type_id in [2, 9, 16, 17]:
            if height < 0.8:
                return 'subtitle'
            if area < 3.0:
                return 'kpi'
            if area < 15.0:
                return 'content'
            return 'main_content'
        
        return 'content'
    
    def _fill_placeholder_smart(self, slide: Any, ph_id: int, ph_info: Dict,
                                section: Any, search_results: Dict, prepared_content: Dict = None) -> Dict:
        """
        Fill a placeholder based on its role and content availability.

        Args:
            slide (Any): The slide object.
            ph_id (int): Placeholder ID.
            ph_info (Dict): Placeholder info.
            section (Any): Section plan.
            search_results (Dict): Search results.
            prepared_content (Dict, optional): Pre-generated content.

        Returns:
            Dict: Log of filling process.
        """
        
        role = ph_info['role']
        area = ph_info['area']
        
        logger.info(f"    Placeholder {ph_id} ({role}, {area:.1f} sq in)")
        
        try:
            placeholder = slide.placeholders[ph_id]
        except KeyError:
            logger.error(f"      ❌ Placeholder {ph_id} not found in slide")
            return {'id': ph_id, 'status': 'not_found'}

        # Try to find icon if content description mentions icon/symbol
        # Or if the role was detected as 'icon' by LLM
        if role == 'icon' or (role == 'content' and area < 1.0):
             # Try to find a keyword for icon
             keyword = section.section_title # Default
             if section.placeholder_specs:
                 for spec in section.placeholder_specs:
                     if spec.placeholder_idx == ph_id:
                         keyword = spec.content_description
                         break

             icon_file = self.icon_selector.select_icon_for_keyword(keyword, self.openai_client)
             if ph_info['type_id'] == 15 or role == 'image':
                 try:
                     from slidedeckai.global_config import GlobalConfig
                     # Get full path for icon using GlobalConfig
                     icon_path = GlobalConfig.ICONS_DIR / icon_file
                     if not icon_path.exists():
                         # Fallback to placeholder if icon not found
                         icon_path = GlobalConfig.ICONS_DIR.parent / "placeholder.png"

                     if icon_path.exists():
                         placeholder.insert_picture(str(icon_path))
                         logger.info(f"      ✓ Icon inserted: {icon_file}")
                         return {'id': ph_id, 'role': role, 'icon': icon_file, 'status': 'filled'}
                 except Exception as e:
                     logger.warning(f"      ⚠️ Failed to insert icon: {e}")
        
        if role == 'subtitle':
            # If pre-generated content exists, use it
            if prepared_content and prepared_content.get(ph_id):
                subtitle_text = prepared_content[ph_id].get('text')
                if subtitle_text:
                    placeholder.text = subtitle_text
                    for para in placeholder.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.name = self.template_properties['default_fonts']['name']
                            try:
                                base_pt = self.template_properties['default_fonts']['size'].pt
                                run.font.size = Pt(base_pt * 0.8)
                            except Exception:
                                run.font.size = Pt(14)
                    return {'id': ph_id, 'role': 'subtitle', 'content': subtitle_text, 'status': 'filled'}
            return self._fill_subtitle(placeholder, ph_id, section, search_results)

        elif role == 'chart':
            # Use pre-generated chart data if available
            pre = prepared_content.get(ph_id) if prepared_content else None
            if pre and pre.get('chart_data'):
                # Insert chart using pre['chart_data']
                try:
                    chart_data_json = pre['chart_data']
                    # reuse _fill_chart logic by temporarily attaching chart_data_json
                    # We'll call a helper that inserts chart_data directly
                    return self._insert_chart_from_data(slide, placeholder, ph_id, chart_data_json)
                except Exception as e:
                    logger.warning(f"Failed to insert pre-generated chart for ph {ph_id}: {e}")
            return self._fill_chart(slide, placeholder, ph_id, section, search_results)

        elif role == 'table':
            pre = prepared_content.get(ph_id) if prepared_content else None
            if pre and pre.get('table_data'):
                try:
                    table_data = pre['table_data']
                    return self._insert_table_from_data(slide, placeholder, ph_id, table_data)
                except Exception as e:
                    logger.warning(f"Failed to insert pre-generated table for ph {ph_id}: {e}")
            return self._fill_table(slide, placeholder, ph_id, section, search_results)
        
        elif role == 'kpi':
            return self._fill_kpi(placeholder, ph_id, ph_info, section, search_results)
        
        elif role in ['content', 'main_content']:
            return self._fill_content(placeholder, ph_id, ph_info, section, search_results)
        
        else:
            logger.warning(f"      ⚠️ Unknown role: {role}")
            return {'id': ph_id, 'role': role, 'status': 'skipped'}
    
    def _fill_subtitle(self, placeholder: Any, ph_id: int, section: Any, search_results: Dict) -> Dict:
        """
        Fill a subtitle placeholder using template fonts safely.

        Args:
            placeholder (Any): The placeholder object.
            ph_id (int): Placeholder ID.
            section (Any): Section plan.
            search_results (Dict): Search results.

        Returns:
            Dict: Log.
        """
        
        if not placeholder.has_text_frame:
            return {'id': ph_id, 'status': 'no_text_frame'}
        
        subtitle = self.content_generator.generate_subtitle(
            section.section_title,
            section.section_purpose,
            list(search_results.values())[:3]
        )
        
        placeholder.text = subtitle
        
        # Apply template font safely
        try:
            default_fonts = self.template_properties.get('default_fonts', {'name': 'Calibri', 'size': Pt(18)})
            font_name = default_fonts.get('name', 'Calibri')
            font_size = default_fonts.get('size', Pt(18))
            
            for para in placeholder.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = font_name
                    try:
                        if hasattr(font_size, 'pt'):
                            base_pt = font_size.pt
                        else:
                            base_pt = 18
                        run.font.size = Pt(base_pt * 0.8)
                    except Exception:
                        run.font.size = Pt(14)  # fallback
        except Exception as e:
            logger.debug(f"Font application failed: {e}")
        
        logger.info(f"      Subtitle: {subtitle}")
        
        return {
            'id': ph_id,
            'role': 'subtitle',
            'content': subtitle,
            'status': 'filled'
        }
    
    def _fill_chart(self, slide: Any, placeholder: Any, ph_id: int, section: Any, search_results: Dict) -> Dict:
        """
        Generate and insert a chart into a placeholder.

        Args:
            slide (Any): The slide.
            placeholder (Any): The placeholder.
            ph_id (int): Placeholder ID.
            section (Any): Section plan.
            search_results (Dict): Search results.

        Returns:
            Dict: Log.
        """
        
        if 'chart' not in section.enforced_content_type:
            return {'id': ph_id, 'status': 'skipped'}
        
        # Gather relevant data
        relevant_facts = []
        for spec in section.placeholder_specs:
            for query in spec.search_queries:
                if query.query in search_results:
                    relevant_facts.extend(search_results[query.query])
        
        # Generate chart data
        chart_data_json = self.content_generator.generate_chart(
            section.section_title,
            section.section_purpose,
            relevant_facts,
            chart_type='column'
        )
        
        logger.info(f"      ✓ Chart: {len(chart_data_json.get('categories', []))} cats")
        
        # Get placeholder position
        x = placeholder.left
        y = placeholder.top
        cx = placeholder.width
        cy = placeholder.height
        
        # Remove placeholder
        try:
            sp = placeholder.element
            sp.getparent().remove(sp)
        except Exception as e:
            logger.warning(f"Could not remove placeholder: {e}")
        
        # Chart type mapping
        chart_type_map = {
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE,
        }
        
        chart_type = chart_data_json.get('type', 'column').lower()
        xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        # Prepare data
        chart_data = CategoryChartData()
        chart_data.categories = chart_data_json.get('categories', [])
        
        series_data = chart_data_json.get('series', [])
        
        if chart_type == 'pie':
            if isinstance(series_data[0], dict):
                chart_data.add_series(
                    series_data[0].get('name', 'Values'), 
                    series_data[0].get('values', [])
                )
            else:
                chart_data.add_series('Values', series_data)
        else:
            for series_item in series_data:
                if isinstance(series_item, dict):
                    chart_data.add_series(
                        series_item.get('name', 'Series'), 
                        series_item.get('values', [])
                    )
                else:
                    chart_data.add_series('Series', series_item)
        
        # INSERT CHART (FIX #3 - COMPLETE)
        try:
            chart_shape = slide.shapes.add_chart(xl_chart_type, x, y, cx, cy, chart_data)
            chart = chart_shape.chart
            chart.has_legend = True
            
            if chart_type != 'pie':
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
            
            if 'title' in chart_data_json:
                chart.has_title = True
                chart.chart_title.text_frame.text = chart_data_json['title']
                # Apply template font to chart title
                try:
                    default_fonts = self.template_properties.get('default_fonts', {'name': 'Calibri', 'size': Pt(18)})
                    font_name = default_fonts.get('name', 'Calibri')
                    font_size = default_fonts.get('size', Pt(18))
                    for para in chart.chart_title.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.name = font_name
                            if hasattr(font_size, 'pt'):
                                run.font.size = font_size
                except Exception:
                    pass
            
            logger.info(f"      ✓ Chart inserted successfully")
        except Exception as e:
            logger.error(f"Failed to insert chart: {e}")
            return {'id': ph_id, 'status': 'failed', 'error': str(e)}
        
        return {
            'id': ph_id,
            'role': 'chart',
            'chart_data': chart_data_json,
            'status': 'filled'
        }

    def _insert_chart_from_data(self, slide: Any, placeholder: Any, ph_id: int, chart_data_json: Dict) -> Dict:
        """
        Insert chart when chart data JSON is already available (no LLM calls).

        Args:
            slide (Any): The slide.
            placeholder (Any): The placeholder.
            ph_id (int): Placeholder ID.
            chart_data_json (Dict): The chart data.

        Returns:
            Dict: Log.
        """
        # mirrored insertion logic from _fill_chart
        try:
            x = placeholder.left
            y = placeholder.top
            cx = placeholder.width
            cy = placeholder.height
            try:
                sp = placeholder.element
                sp.getparent().remove(sp)
            except Exception:
                pass

            chart_type_map = {
                'bar': XL_CHART_TYPE.BAR_CLUSTERED,
                'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'line': XL_CHART_TYPE.LINE,
                'pie': XL_CHART_TYPE.PIE,
            }

            chart_type = chart_data_json.get('type', 'column').lower()
            xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

            chart_data = CategoryChartData()
            chart_data.categories = chart_data_json.get('categories', [])
            series_data = chart_data_json.get('series', [])
            if chart_type == 'pie':
                if isinstance(series_data[0], dict):
                    chart_data.add_series(series_data[0].get('name', 'Values'), series_data[0].get('values', []))
                else:
                    chart_data.add_series('Values', series_data)
            else:
                for series_item in series_data:
                    if isinstance(series_item, dict):
                        chart_data.add_series(series_item.get('name', 'Series'), series_item.get('values', []))
                    else:
                        chart_data.add_series('Series', series_item)

            chart_shape = slide.shapes.add_chart(xl_chart_type, x, y, cx, cy, chart_data)
            chart = chart_shape.chart
            chart.has_legend = True
            if chart_type != 'pie':
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False

            if 'title' in chart_data_json:
                chart.has_title = True
                chart.chart_title.text_frame.text = chart_data_json['title']

            return {'id': ph_id, 'role': 'chart', 'chart_data': chart_data_json, 'status': 'filled'}
        except Exception as e:
            logger.error(f"_insert_chart_from_data failed: {e}")
            return {'id': ph_id, 'status': 'failed', 'error': str(e)}
    
    def _fill_table(self, slide: Any, placeholder: Any, ph_id: int, section: Any, search_results: Dict) -> Dict:
        """
        Generate and insert a table into a placeholder.

        Args:
            slide (Any): The slide.
            placeholder (Any): The placeholder.
            ph_id (int): Placeholder ID.
            section (Any): Section plan.
            search_results (Dict): Search results.

        Returns:
            Dict: Log.
        """
        
        if 'table' not in section.enforced_content_type:
            return {'id': ph_id, 'status': 'skipped'}
        
        relevant_facts = []
        for spec in section.placeholder_specs:
            for query in spec.search_queries:
                if query.query in search_results:
                    relevant_facts.extend(search_results[query.query])
        
        table_data = self.content_generator.generate_table(
            section.section_title,
            section.section_purpose,
            relevant_facts
        )
        
        headers = table_data.get('headers', [])
        rows = table_data.get('rows', [])
        
        if not headers or not rows:
            return {'id': ph_id, 'status': 'no_data'}
        
        logger.info(f"      ✓ Table: {len(headers)} cols")
        
        # Get placeholder dimensions
        left = placeholder.left
        top = placeholder.top
        width = placeholder.width
        height = placeholder.height
        
        # Remove placeholder
        try:
            sp = placeholder.element
            sp.getparent().remove(sp)
        except Exception as e:
            logger.warning(f"Could not remove placeholder: {e}")
        
        # INSERT TABLE (FIX #3 - COMPLETE)
        try:
            from pptx.util import Inches
            
            table_shape = slide.shapes.add_table(
                len(rows) + 1, 
                len(headers), 
                left, top, width, height
            )
            table = table_shape.table
            
            # Smart column widths
            col_max_lengths = [len(str(headers[col])) for col in range(len(headers))]
            for row_data in rows:
                for col_idx in range(len(headers)):
                    if col_idx < len(row_data):
                        col_max_lengths[col_idx] = max(
                            col_max_lengths[col_idx], 
                            len(str(row_data[col_idx]))
                        )
            
            total_length = sum(col_max_lengths)
            if total_length > 0:
                for col_idx in range(len(headers)):
                    proportion = max(col_max_lengths[col_idx] / total_length, 0.1)
                    table.columns[col_idx].width = int(width * proportion)
            
            # Headers
            for col_idx, header_text in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header_text)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        default_fonts = self.template_properties.get('default_fonts', {'name': 'Calibri', 'size': Pt(18)})
                        run.font.name = default_fonts.get('name', 'Calibri')
            
            # Rows
            for row_idx, row_data in enumerate(rows, start=1):
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx < len(headers):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(cell_text)
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                default_fonts = self.template_properties.get('default_fonts', {'name': 'Calibri', 'size': Pt(18)})
                                run.font.name = default_fonts.get('name', 'Calibri')
            
            logger.info(f"      ✓ Table inserted successfully")
        except Exception as e:
            logger.error(f"Failed to insert table: {e}")
            return {'id': ph_id, 'status': 'failed', 'error': str(e)}
        
        return {
            'id': ph_id,
            'role': 'table',
            'table_data': table_data,
            'status': 'filled'
        }

    def _insert_table_from_data(self, slide: Any, placeholder: Any, ph_id: int, table_data: Dict) -> Dict:
        """
        Insert table when table data is already available.

        Args:
            slide (Any): The slide.
            placeholder (Any): The placeholder.
            ph_id (int): Placeholder ID.
            table_data (Dict): Table data.

        Returns:
            Dict: Log.
        """
        headers = table_data.get('headers', [])
        rows = table_data.get('rows', [])
        if not headers or not rows:
            return {'id': ph_id, 'status': 'no_data'}

        left = placeholder.left
        top = placeholder.top
        width = placeholder.width
        height = placeholder.height

        try:
            sp = placeholder.element
            sp.getparent().remove(sp)
        except Exception:
            pass

        try:
            table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
            table = table_shape.table
            col_max_lengths = [len(str(headers[col])) for col in range(len(headers))]
            for row_data in rows:
                for col_idx in range(len(headers)):
                    if col_idx < len(row_data):
                        col_max_lengths[col_idx] = max(col_max_lengths[col_idx], len(str(row_data[col_idx])))

            total_length = sum(col_max_lengths)
            if total_length > 0:
                for col_idx in range(len(headers)):
                    proportion = max(col_max_lengths[col_idx] / total_length, 0.1)
                    table.columns[col_idx].width = int(width * proportion)

            for col_idx, header_text in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header_text)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        try:
                            run.font.name = self.template_properties['default_fonts']['name']
                        except Exception:
                            pass

            for row_idx, row_data in enumerate(rows, start=1):
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx < len(headers):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(cell_text)
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                try:
                                    run.font.name = self.template_properties['default_fonts']['name']
                                except Exception:
                                    pass

            return {'id': ph_id, 'role': 'table', 'table_data': table_data, 'status': 'filled'}
        except Exception as e:
            logger.error(f"_insert_table_from_data failed: {e}")
            return {'id': ph_id, 'status': 'failed', 'error': str(e)}
    
    def _fill_kpi(self, placeholder: Any, ph_id: int, ph_info: Dict,
                  section: Any, search_results: Dict) -> Dict:
        """
        Fill a KPI placeholder with theme colors.

        Args:
            placeholder (Any): The placeholder.
            ph_id (int): Placeholder ID.
            ph_info (Dict): Placeholder info.
            section (Any): Section plan.
            search_results (Dict): Search results.

        Returns:
            Dict: Log.
        """
        
        if not placeholder.has_text_frame:
            return {'id': ph_id, 'status': 'no_text_frame'}
        
        relevant_facts = []
        for spec in section.placeholder_specs:
            for query in spec.search_queries:
                if query.query in search_results:
                    relevant_facts.extend(search_results[query.query][:1])
        
        if not relevant_facts:
            relevant_facts = [f"KPI for {section.section_title}"]
        
        kpi_data = self.content_generator.generate_kpi(
            section.section_title,
            relevant_facts[0]
        )
        
        text_frame = placeholder.text_frame
        text_frame.clear()
        
        # Value (big)
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = kpi_data['value']
        run.font.size = self._calculate_font_size_from_area(ph_info['area'], 'large')
        run.font.bold = True
        default_fonts = self.template_properties.get('default_fonts', {'name': 'Calibri', 'size': Pt(18)})
        run.font.name = default_fonts.get('name', 'Calibri')
        # FIX #4: Use theme color
        # Prefer actual theme color RGB if available
        try:
            accent = self.template_properties.get('theme_colors', {}).get('accent1')
            if hasattr(accent, 'rgb'):
                run.font.color.rgb = accent.rgb
            elif isinstance(accent, bytes) or isinstance(accent, str):
                # try parsing string
                try:
                    run.font.color.rgb = RGBColor.from_string(accent)
                except Exception:
                    run.font.color.theme_color = 4
            else:
                run.font.color.theme_color = 4
        except Exception:
            run.font.color.theme_color = 4
        
        # Label (small)
        p = text_frame.add_paragraph()
        run = p.add_run()
        run.text = kpi_data['label']
        run.font.size = self._calculate_font_size_from_area(ph_info['area'], 'small')
        run.font.name = default_fonts.get('name', 'Calibri')
        
        logger.info(f"      ✓ KPI: {kpi_data['label']}")
        
        return {
            'id': ph_id,
            'role': 'kpi',
            'kpi_data': kpi_data,
            'status': 'filled'
        }
    
    def _fill_content(self, placeholder: Any, ph_id: int, ph_info: Dict,
                      section: Any, search_results: Dict) -> Dict:
        """
        Fill a general content placeholder using template fonts.

        Args:
            placeholder (Any): The placeholder.
            ph_id (int): Placeholder ID.
            ph_info (Dict): Placeholder info.
            section (Any): Section plan.
            search_results (Dict): Search results.

        Returns:
            Dict: Log.
        """
        
        if not placeholder.has_text_frame:
            return {'id': ph_id, 'status': 'no_text_frame'}
        
        relevant_facts = []
        for spec in section.placeholder_specs:
            for query in spec.search_queries:
                if query.query in search_results:
                    relevant_facts.extend(search_results[query.query])
        
        max_bullets = self._calculate_max_bullets(ph_info['area'])
        
        bullets = self.content_generator.generate_bullets(
            section.section_title,
            section.section_purpose,
            relevant_facts,
            max_bullets=max_bullets
        )
        
        text_frame = placeholder.text_frame
        text_frame.clear()
        
        for idx, bullet in enumerate(bullets):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet
            p.level = 0
            
            # FIX #4: Apply template font
            for run in p.runs:
                default_fonts = self.template_properties.get('default_fonts', {'name': 'Calibri', 'size': Pt(18)})
                run.font.name = default_fonts.get('name', 'Calibri')
                try:
                    font_size = default_fonts.get('size', Pt(18))
                    if hasattr(font_size, 'pt'):
                        run.font.size = font_size
                except Exception:
                    pass
        
        logger.info(f"      ✓ {len(bullets)} bullets")
        
        return {
            'id': ph_id,
            'role': 'content',
            'bullets': bullets,
            'status': 'filled'
        }
    
    def _calculate_max_bullets(self, area: float) -> int:
        """
        Calculate maximum number of bullets based on placeholder area.

        Args:
            area (float): Area in square inches.

        Returns:
            int: Max bullets.
        """
        if area < 3:
            return 3
        elif area < 10:
            return 5
        elif area < 20:
            return 7
        else:
            return 10
    
    def _calculate_font_size_from_area(self, area: float, size_type: str) -> Pt:
        """
        Calculate appropriate font size based on area and type.

        Args:
            area (float): Area in square inches.
            size_type (str): 'large' or 'small'.

        Returns:
            Pt: Font size.
        """
        from pptx.util import Pt
        
        base_size = self.template_properties['default_fonts']['size'].pt
        
        if size_type == 'large':
            if area < 2:
                return Pt(base_size * 1.3)
            elif area < 5:
                return Pt(base_size * 1.8)
            else:
                return Pt(base_size * 2.5)
        else:  # small
            if area < 2:
                return Pt(base_size * 0.6)
            elif area < 5:
                return Pt(base_size * 0.7)
            else:
                return Pt(base_size * 0.8)
    
    def _batch_validate_placeholder_roles(self, section: Any, placeholder_map: Dict) -> Dict:
        """
        Batch-validate placeholder roles using LLM.

        Args:
            section (Any): Section plan.
            placeholder_map (Dict): Placeholder map.

        Returns:
            Dict: Mapping of {ph_id: valid_role}.
        """
        roles = ['subtitle', 'chart', 'table', 'kpi', 'content', 'main_content', 'image', 'icon']
        items = []
        for pid, info in placeholder_map.items():
            try:
                pid_int = int(pid)
            except Exception:
                pid_int = pid
            items.append({
                'id': pid_int,
                'type': info.get('type'),
                'area': round(float(info.get('area', 0)), 2),
                'inferred_role': info.get('role')
            })

        prompt = {
            'section_title': getattr(section, 'section_title', ''),
            'section_purpose': getattr(section, 'section_purpose', ''),
            'placeholders': items,
            'allowed_roles': roles
        }

        instruction = (
            "Given section context and placeholders, return JSON mapping placeholder ids to best role from allowed_roles. Return ONLY JSON."
        )

        messages = [
            {"role": "system", "content": "You are a concise classifier. Return only valid JSON."},
            {"role": "user", "content": instruction + "\n\n" + json.dumps(prompt)}
        ]

        try:
            resp = self.content_generator.client.chat.completions.create(
                model=self.content_generator.model,
                messages=messages,
                temperature=0.0,
                max_tokens=400
            )

            text = resp.choices[0].message.content.strip()
            try:
                parsed = json.loads(text)
            except Exception:
                import re
                m = re.search(r"\{[\s\S]*\}", text)
                if m:
                    parsed = json.loads(m.group(0))
                else:
                    raise

            result = {}
            for k, v in parsed.items():
                try:
                    idx = int(k)
                except Exception:
                    try:
                        idx = int(float(k))
                    except Exception:
                        continue
                role = str(v).lower()
                if role not in roles:
                    for r in roles:
                        if r in role:
                            role = r
                            break
                    else:
                        role = placeholder_map.get(idx, {}).get('role')
                result[idx] = role

            return result
        except Exception as e:
            logger.debug(f"Batch role validation failed: {e}")
            return {int(pid): info.get('role') for pid, info in placeholder_map.items()}

    def _execute_mock_plan(self, plan: Any, output_path: pathlib.Path) -> pathlib.Path:
        """
        Execute a plan in demo mode using mock data.

        Args:
            plan (Any): The plan.
            output_path (pathlib.Path): Output path.

        Returns:
            pathlib.Path: Path to saved file.
        """

        # Add Title Slide
        self._add_title_slide(plan.query)

        for section in plan.sections:
            layout_idx = section.layout_idx
            layout = self.presentation.slide_layouts[layout_idx]
            slide = self.presentation.slides.add_slide(layout)

            # Title
            if slide.shapes.title:
                slide.shapes.title.text = section.section_title

            # Mock content for placeholders
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 0: continue

                # Simple fallback filling
                if shape.has_text_frame:
                    shape.text = f"Demo Content for {section.section_purpose}\n- Mock Point 1\n- Mock Point 2"

        # Add Thank You
        self._add_thank_you_slide()

        self.presentation.save(output_path)

        # Save mock log
        log_path = str(output_path).replace('.pptx', '.execution.json')
        with open(log_path, 'w') as f:
            json.dump([{'slide': 1, 'status': 'demo_success'}], f)

        return output_path

    def _get_placeholder_type_name(self, type_id: int) -> str:
        """
        Get the string name for a placeholder type ID.

        Args:
            type_id (int): Type ID.

        Returns:
            str: Type name.
        """
        TYPES = {
            1: 'TITLE', 2: 'BODY', 3: 'CENTER_TITLE', 4: 'SUBTITLE',
            5: 'DATE', 6: 'SLIDE_NUMBER', 7: 'FOOTER', 8: 'HEADER',
            9: 'OBJECT', 10: 'CHART', 11: 'TABLE', 12: 'CLIP_ART',
            13: 'ORG_CHART', 14: 'MEDIA', 15: 'PICTURE',
            16: 'VERTICAL_BODY', 17: 'VERTICAL_OBJECT', 18: 'VERTICAL_TITLE',
        }
        return TYPES.get(type_id, f'UNKNOWN_{type_id}')
